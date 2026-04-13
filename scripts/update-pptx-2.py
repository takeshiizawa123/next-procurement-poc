"""
PPTXファイル3本を2026-04-13マニュアル整備に合わせて更新
追加内容: 監査ログ、DLQ、DBバックアップ、全文検索、Slack AI、科目修正学習ループ
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def replace_text_in_slide(slide, replacements):
    """スライド内のすべてのシェイプでテキストを置換"""
    changed = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            changed = True
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for old, new in replacements.items():
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                        changed = True
    return changed


def add_slide_with_content(prs, title, bullets, layout_idx=1):
    """新規スライドを追加（タイトル + 箇条書き）"""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # タイトル
    if slide.shapes.title:
        slide.shapes.title.text = title
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.bold = True
    # 本文
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # body placeholder
            body = shape
            break
    if body and body.has_text_frame:
        body.text_frame.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = body.text_frame.paragraphs[0]
            else:
                p = body.text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.space_after = Pt(6)
    return slide


# ========================================================
# 共通置換
# ========================================================
COMMON_REPLACEMENTS_2 = {
    "15テーブル": "18テーブル",
    "v2.0 — 2026-04-12（Supabase版）": "v3.0 — 2026-04-13",
    "v3.0（Supabase版）": "v3.1（監査ログ・AI検索対応）",
    "最終更新: 2026-04-12": "最終更新: 2026-04-13",
    "2026-04-12": "2026-04-13",
}

# ========================================================
# 1. operational-guide.pptx
# ========================================================
print("=== Updating operational-guide.pptx ===")
prs1 = Presentation("docs/operational-guide.pptx")

OPS_REPLACEMENTS = {
    "Supabase版 全機能実装完了": "Supabase版 全機能+障害対策基盤 実装完了",
    "Stage2/3仕訳自動化済み → 本番展開準備中": "監査ログ・DLQ・バックアップ・AI検索・学習ループ完備",
}

for slide in prs1.slides:
    replace_text_in_slide(slide, COMMON_REPLACEMENTS_2)
    replace_text_in_slide(slide, OPS_REPLACEMENTS)

# 新スライド: 障害対策基盤
add_slide_with_content(prs1, "障害対策基盤（§13）", [
    "監査ログ（audit_log）: ステータス変更時に自動記録",
    "リトライ+DLQ: 外部API失敗→指数バックオフ4回→OPS通知",
    "日次DBバックアップ: JST 03:00 → Google Drive永久保持",
    "科目修正学習ループ: 仕訳画面の修正→RAGコンテキスト注入",
    "障害復旧手順書: docs/disaster-recovery.md",
])

# 新スライド: Slack AI + 全文検索
add_slide_with_content(prs1, "Slack AI・全文検索（§14-15）", [
    "Slack AIアシスタント（/ask）: Claude Haiku + RAG応答",
    "  → 購買・出張に関する質問をSlackから即座に回答",
    "全文検索: マイページ検索バー（pg_trgm + GIN）",
    "  → 品目名・仕入先名・備考を部分一致検索",
    "環境変数: ANTHROPIC_API_KEY（AI用）",
])

# 新スライド: Cron一覧（更新版）
add_slide_with_content(prs1, "定期実行一覧（8 cron jobs）", [
    "cache-warm: 4分毎 — Redisキャッシュ先読み",
    "daily-summary: 09:00 — OPSチャンネルに日次サマリ",
    "voucher-reminder: 10:00 — 証憑催促（Day1/3/7）",
    "weekly-reminder: 月曜09:00 — 承認待ち週次まとめ",
    "card-reconciliation: 月曜11:00 — カード明細照合",
    "daily-variance: 12:00 — 金額乖離検知",
    "trip-controls: 毎月1日10:00 — 出張統制レポート",
    "db-backup: 03:00 — DBバックアップ→GDrive永久保持",
])

prs1.save("docs/operational-guide.pptx")
print("  Saved operational-guide.pptx")

# ========================================================
# 2. user-manual.pptx
# ========================================================
print("=== Updating user-manual.pptx ===")
prs2 = Presentation("docs/user-manual.pptx")

USER_REPLACEMENTS_2 = {
    "v1.3": "v1.4",
}

for slide in prs2.slides:
    replace_text_in_slide(slide, COMMON_REPLACEMENTS_2)
    replace_text_in_slide(slide, USER_REPLACEMENTS_2)

# 新スライド: 検索とAI
add_slide_with_content(prs2, "新機能: 検索・AIアシスタント", [
    "マイページ全文検索（/purchase/my）",
    "  → 検索バーにキーワード入力で過去の申請を検索",
    "  → 品目名・仕入先名・備考を部分一致",
    "",
    "Slack AIアシスタント（/ask）",
    "  → Slackから購買・出張に関する質問ができる",
    "  → 例: /ask 先月のモニター購入は？",
    "  → 購買データに基づいたAI回答が即座に返る",
])

prs2.save("docs/user-manual.pptx")
print("  Saved user-manual.pptx")

# ========================================================
# 3. workflow-design-b-route.pptx
# ========================================================
print("=== Updating workflow-design-b-route.pptx ===")
prs3 = Presentation("docs/workflow-design-b-route.pptx")

WORKFLOW_REPLACEMENTS_2 = {
    "本番展開\n管理本部→他部門の段階展開（準備中）":
        "本番展開（準備中）\n監査ログ・DLQ・バックアップ・AI検索完備\n管理本部→他部門の段階展開",
    "Stage 2/3自動化 ✓完了\nカード照合確定・引落消込で仕訳自動作成":
        "Stage 2/3自動化 ✓完了\nカード照合確定・引落消込で仕訳自動作成\n+ 科目修正学習ループ",
}

for slide in prs3.slides:
    replace_text_in_slide(slide, COMMON_REPLACEMENTS_2)
    replace_text_in_slide(slide, WORKFLOW_REPLACEMENTS_2)

# 新スライド: 障害対策・運用基盤
add_slide_with_content(prs3, "障害対策・運用基盤（実装済み）", [
    "監査ログ: 全ステータス変更を audit_log テーブルに自動記録",
    "リトライ+DLQ: 指数バックオフ(4回) → dead_letter_queue → OPS通知",
    "日次バックアップ: Supabase全テーブル → GDrive JSON（永久保持）",
    "勘定科目学習: account_corrections → RAGコンテキスト注入",
    "全文検索: pg_trgm + GINインデックス（品目名/仕入先/備考）",
    "Slack AI: /askコマンド → Claude Haiku + 購買データRAG",
    "出張統制: 月次cron → 差異検知・未申請・重複・部門別コスト",
])

prs3.save("docs/workflow-design-b-route.pptx")
print("  Saved workflow-design-b-route.pptx")

print("\n=== All 3 PPTX files updated for 2026-04-13 ===")

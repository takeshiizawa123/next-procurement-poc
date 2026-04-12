"""
PPTXファイル3本をSupabase版に更新するスクリプト
テキスト置換 + スライド内容更新
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from copy import deepcopy


def replace_text_in_shape(shape, replacements):
    """シェイプ内のテキストを置換（書式保持）"""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    changed = True
    return changed


def replace_text_in_slide(slide, replacements):
    """スライド内のすべてのシェイプでテキストを置換"""
    for shape in slide.shapes:
        replace_text_in_shape(shape, replacements)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for old, new in replacements.items():
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)


def set_slide_text(slide, texts_by_position):
    """スライド内のテキストフレームを位置順に更新"""
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    text_shapes.sort(key=lambda s: (s.top, s.left))
    for idx, new_text in texts_by_position.items():
        if idx < len(text_shapes):
            shape = text_shapes[idx]
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            if shape.text_frame.paragraphs:
                shape.text_frame.paragraphs[0].runs[0].text = new_text if shape.text_frame.paragraphs[0].runs else new_text


# ========================================================
# 共通置換（全3ファイル共通）
# ========================================================
COMMON_REPLACEMENTS = {
    "GAS": "Supabase",
    "購買台帳(GAS)": "Supabase Postgres",
    "購買台帳（GAS）": "Supabase Postgres",
    "Google Apps Script": "Next.js + Vercel",
    "スプレッドシート": "PostgreSQL",
    "card_last4": "office_member_id",
    "カード下4桁": "office_member_id",
    "v2": "v3.0（Supabase版）",
    "v1.0 — 2026-03-28": "v2.0 — 2026-04-12（Supabase版）",
    "2026-03-28  v2": "2026-04-12  v3.0（Supabase版）",
}

# ========================================================
# 1. operational-guide.pptx
# ========================================================
print("=== Updating operational-guide.pptx ===")
prs1 = Presentation("docs/operational-guide.pptx")

# Slide 21: /trip出張申請フロー → 予約完了申請に変更
TRIP_REPLACEMENTS = {
    "/trip 出張申請フロー": "出張予約完了申請フロー（事後承認）",
    "/trip申請": "各サービスで予約・MFカード決済",
    "#出張CH投稿": "/trip/new で予約完了申請",
    "じゃらん/ANA/JAL予約": "#出張CH投稿+部門長事後承認",
    "MFカード決済": "AIアシスタントで自動入力可",
    "CSV取込(月次)": "カード明細自動照合",
    "MF経費→仕訳": "MF経費→MF会計Plus仕訳",
}

for slide in prs1.slides:
    replace_text_in_slide(slide, COMMON_REPLACEMENTS)
    replace_text_in_slide(slide, TRIP_REPLACEMENTS)

# Slide 8: パターンC に /expense/new 追記
slide8 = prs1.slides[7]
for shape in slide8.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "MF経費経由で給与精算" in run.text:
                    run.text = run.text + "\n専用ページ: /expense/new から簡易申請可"

# Slide 33: Ready to Launch → 更新
slide33 = prs1.slides[32]
for shape in slide33.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "全機能完成" in run.text:
                    run.text = "Supabase版 全機能実装完了"
                if "旧WFから一括移行" in run.text:
                    run.text = "Stage2/3仕訳自動化済み → 本番展開準備中"

prs1.save("docs/operational-guide.pptx")
print("  Saved operational-guide.pptx")

# ========================================================
# 2. user-manual.pptx
# ========================================================
print("=== Updating user-manual.pptx ===")
prs2 = Presentation("docs/user-manual.pptx")

USER_MANUAL_REPLACEMENTS = {
    "/trip を入力": "/trip/new（Webフォーム）にアクセス",
    "モーダルフォームが開きます": "AIアシスタント付きフォームが開きます",
    "モーダルに入力": "フォームに入力（AI自動入力可）",
    "行き先・日程・目的交通手段・概算額・宿泊先HubSpot案件番号（任意）": "行き先・日程・目的\n交通費（複数行）・宿泊費（複数行）\nPJコード・HubSpot案件番号",
    "#出張チャンネルに日当自動計算付きで投稿されます": "#出張チャンネルに実額+日当自動計算で投稿",
    "各交通機関・宿泊を予約（次スライド参照）": "事後承認: 部門長が実額確認→[承認]",
    "すべてMFビジネスカードで支払い": "カード明細で自動照合→仕訳",
    "カード明細→MF経費→管理本部が仕訳": "予約完了後に申請 → 事後承認",
    "/trip → モーダル": "/trip/new + AIアシスタント",
    "出張申請 | /trip → モーダル | 申請者": "出張予約完了申請 | /trip/new + AIアシスタント | 申請者",
    "EX予約の場合はICカードで乗車": "EX予約はICカード乗車",
    "タクシー | 現地で利用 | MFカード | 領収書不要（カード明細で確認）": "タイムズカー | タイムズカーアプリ | MFカード | 法人カードを事前登録",
    "レンタカー | 各社Webサイト | MFカード | 利用交通手段欄に「レンタカー」と記入": "レンタカー | トヨタレンタカー等 | MFカード | 法人カードを事前登録",
    "楽天トラベルRacco も今後利用可能になります": "楽天トラベルRacco も利用可能（予約リンク生成対応済み）",
    "立替が発生した場合 → /purchase「購入済」で申請": "立替が発生した場合 → /expense/new で立替精算申請",
}

for slide in prs2.slides:
    replace_text_in_slide(slide, COMMON_REPLACEMENTS)
    replace_text_in_slide(slide, USER_MANUAL_REPLACEMENTS)

prs2.save("docs/user-manual.pptx")
print("  Saved user-manual.pptx")

# ========================================================
# 3. workflow-design-b-route.pptx
# ========================================================
print("=== Updating workflow-design-b-route.pptx ===")
prs3 = Presentation("docs/workflow-design-b-route.pptx")

WORKFLOW_REPLACEMENTS = {
    "購買管理システム | 申請・承認・照合・仕訳作成のすべての司令塔（自社開発・Vercel）":
        "購買管理システム | 申請・承認・照合・仕訳作成の司令塔（Next.js + Vercel + Supabase Postgres）",
    "購買管理システム（司令塔）": "購買管理システム（司令塔 / Supabase Postgres）",
    "申請・承認・予測テーブル・照合エンジン・仕訳作成": "申請・承認・予測テーブル(predicted_transactions)・照合エンジン(card-matcher-v2)・仕訳作成",
    "card_last4 × 金額 × 予想利用日": "office_member_id × 金額 × 予想利用日",
    "照合キー: office_member_id (従業員) × 金額 × 日付": "照合キー: office_member_id × 金額 × 日付（±7日）",
    "□ MFビジネスカード → MF経費 の自動連携設定方法": "検証中: MFビジネスカード → MF経費 の自動連携設定",
    "□ カード明細の反映タイミング（即時 or 翌営業日）": "検証中: カード明細の反映タイミング",
    "□ 同日同額の複数決済の区別方法": "検証中: 同日同額の複数決済の区別方法",
    "□ 返金・キャンセル時の明細形式": "検証中: 返金・キャンセル時の明細形式",
    "□ MF会計Plus の自動仕訳ルールをOFFにできるか": "検証中: MF会計Plus の自動仕訳ルールOFF確認",
    "□ 二重計上防止のための一意キー設計": "実装済み: PO番号による仕訳冪等性チェック",
    "□ 立替精算のフロー（MF経費を通すか直接仕訳か）": "実装済み: /expense/new → MF経費経由精算",
    "✓ MF経費APIで office_member_id が取得可能（実機確認済み）": "✓ office_member_id ベース照合(card-matcher-v2)実装済み",
    # Phase status updates
    "基盤構築\n従業員マスタ拡張・予測テーブル統一・出張申請Web化": "基盤構築 ✓完了\nSupabase Postgres + Drizzle ORM",
    "MF経費連携\nfetchCardStatements改修・office_member_id ベース照合": "MF経費連携 ✓完了\ncard-matcher-v2 実装済み",
    "出張承認\n/trip 承認フロー追加・部門長DMボタン": "出張承認 ✓完了\n予約完了申請(事後承認)+AIアシスタント",
    "立替Web化\n/expense/new 新規・OCR連携": "立替Web化 ✓完了\n/expense/new 実装済み",
    "Stage 2自動化\n購買管理がMF会計Plus APIでStage 2仕訳作成": "Stage 2/3自動化 ✓完了\nカード照合確定・引落消込で仕訳自動作成",
    "本番展開\n管理本部→他部門の段階展開・運用マニュアル整備": "本番展開\n管理本部→他部門の段階展開（準備中）",
}

for slide in prs3.slides:
    replace_text_in_slide(slide, COMMON_REPLACEMENTS)
    replace_text_in_slide(slide, WORKFLOW_REPLACEMENTS)

prs3.save("docs/workflow-design-b-route.pptx")
print("  Saved workflow-design-b-route.pptx")

print("\n=== All 3 PPTX files updated ===")

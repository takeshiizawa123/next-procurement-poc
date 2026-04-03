"""Generate travel service recommendation PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x2D, 0x3A, 0x4A)
GRAY_TEXT = RGBColor(0x6B, 0x7B, 0x8D)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ROW_ALT = RGBColor(0xEB, 0xF0, 0xF5)
TABLE_BORDER = RGBColor(0xD5, 0xDB, 0xE1)


def add_header_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_text(slide, text, left, top, width, height, size=16, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullets(slide, items, left, top, width, height, size=15, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf


def styled_table(slide, data, left, top, width, height):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(12)
            p.font.bold = (i == 0)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                p.font.color.rgb = WHITE
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
                p.font.color.rgb = DARK_TEXT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                p.font.color.rgb = DARK_TEXT
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return table


def add_card(slide, x, y, w, h, title, items, color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = TABLE_BORDER
    card.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, title, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
             size=16, color=PRIMARY, bold=True)
    add_bullets(slide, items, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), h - Inches(0.7),
                size=13, color=DARK_TEXT)


def add_phase_card(slide, x, y, w, h, title, items, color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    add_text(slide, title, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.4),
             size=16, color=WHITE, bold=True)
    add_bullets(slide, items, x + Inches(0.15), y + Inches(0.75), w - Inches(0.3), h - Inches(1.0),
                size=13, color=DARK_TEXT)


# =====================================================
# Slide 1: Title
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY
bg.line.fill.background()

stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(2.5))
stripe.fill.solid()
stripe.fill.fore_color.rgb = RGBColor(0x15, 0x2E, 0x4A)
stripe.line.fill.background()

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.0), Inches(0.08), Inches(2.0))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_text(slide, "出張手配サービス導入 推奨案", Inches(1.5), Inches(3.0), Inches(10), Inches(0.9),
         size=40, color=WHITE, bold=True)
add_text(slide, "会社メール個人アカウント + MFビジネスカード統一運用",
         Inches(1.5), Inches(3.8), Inches(10), Inches(0.6), size=22, color=ACCENT)
add_text(slide, "フューチャースタンダード 管理本部", Inches(1.5), Inches(5.5), Inches(5), Inches(0.4),
         size=16, color=GRAY_TEXT)
add_text(slide, "2026年4月", Inches(1.5), Inches(5.9), Inches(5), Inches(0.4),
         size=16, color=GRAY_TEXT)

# =====================================================
# Slide 2: Executive Summary
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "エグゼクティブサマリー", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.2))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
box.line.color.rgb = ACCENT
box.line.width = Pt(2)

add_text(slide, "推奨: 全サービス「会社メール個人アカウント + MFカード」で統一（法人契約なし）",
         Inches(1.2), Inches(1.7), Inches(11), Inches(0.5), size=20, color=PRIMARY, bold=True)
add_text(slide, "即日導入・運用統一・MF自動取込の3点を最適化。必要に応じて法人契約を段階的に追加",
         Inches(1.2), Inches(2.2), Inches(11), Inches(0.4), size=16, color=DARK_TEXT)

cw = Inches(3.6)
ch = Inches(3.5)
cy = Inches(3.3)
add_card(slide, Inches(0.8), cy, cw, ch, "即日導入", [
    "法人契約・審査が一切不要",
    "会社メールでアカウント作成するだけ",
    "全サービス同日に利用開始可能",
    "→ 最短で明日から運用できる",
], GREEN)
add_card(slide, Inches(4.85), cy, cw, ch, "運用統一", [
    "全サービス同一パターンで管理が簡単",
    "入社時: アカウント作成+MFカード登録",
    "退職時: MFカード解除で完結",
    "→ 管理本部の運用負荷を最小化",
], ACCENT)
add_card(slide, Inches(8.9), cy, cw, ch, "MF自動取込", [
    "全出張費がMFカード明細で自動取込",
    "MF会計Plusに集約 → 立替精算なし",
    "部門別・PJ別のコスト可視化",
    "→ 法人契約追加で更に詳細化も可能",
], ORANGE)

# =====================================================
# Slide 3: Recommended Configuration (案A: 個人アカウント統一)
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "推奨案: 全サービス個人アカウント統一", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

styled_table(slide, [
    ["手配対象", "サービス", "アカウント", "決済"],
    ["宿泊", "じゃらん / 楽天トラベル（一般）", "会社メールで個人アカウント", "MFカード"],
    ["航空券（ANA）", "ANA公式サイト", "会社メールで個人アカウント", "MFカード"],
    ["航空券（JAL）", "JAL公式サイト", "会社メールで個人アカウント", "MFカード"],
    ["東海道・山陽新幹線", "スマートEX", "会社メールで個人アカウント", "MFカード"],
    ["東北・北陸新幹線等", "えきねっと", "会社メールで個人アカウント", "MFカード"],
    ["レンタカー", "各社サイト", "会社メールで個人アカウント", "MFカード"],
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.2))

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.fill.background()

add_text(slide, "前提: 全社員にMFビジネスカード配布", Inches(1.2), Inches(5.15), Inches(5), Inches(0.4),
         size=18, color=PRIMARY, bold=True)

styled_table(slide, [
    ["カード種別", "用途", "備考"],
    ["物理カード", "実店舗・緊急購入・出張先", "社員が常時携帯"],
    ["バーチャルカード", "ECサイト・オンライン予約", "用途別に複数発行も可"],
], Inches(1.2), Inches(5.55), Inches(5.0), Inches(1.1))

add_bullets(slide, [
    "全サービス同一パターン: 会社メールで個人アカウント + MFカード紐付け",
    "法人契約ゼロ — 審査不要、即日導入可能",
    "MFカード明細でMF会計Plusに自動取込 → 立替精算なし",
    "※ 現在一部社員のみ配布済 → 全社員への追加発行が必要",
], Inches(6.6), Inches(5.15), Inches(5.8), Inches(1.6), size=13, color=DARK_TEXT)

# =====================================================
# Slide 4: 全個人 vs 宿泊のみ法人契約 比較
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "比較: 全て個人アカウント vs 宿泊のみ法人契約を追加", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=26, color=WHITE, bold=True)

add_text(slide, "※ 航空券・新幹線・レンタカーはどちらの案も「個人アカウント + MFカード」で共通",
         Inches(0.8), Inches(1.15), Inches(11), Inches(0.3), size=13, color=GRAY_TEXT)

styled_table(slide, [
    ["比較項目", "案A: 全て個人アカウント", "案B: 宿泊のみ法人契約（JCS or Racco）"],
    ["導入スピード", "○ 即日、審査なし", "△ JCS: リクルート審査 / Racco: 与信審査2週間"],
    ["PJ番号の紐付け", "○ /trip申請時にHubSpot案件番号で紐付け済", "○ /trip + JCSカスタム項目（二重管理）"],
    ["宿泊の経費把握", "△ MFカード明細のみ（宿名なし）", "○ 管理画面CSV（宿名・日程・金額）"],
    ["上限金額の制御", "× なし（社員の自律に依存）", "○ JCS: エリア毎に上限金額設定可"],
    ["API連携（将来）", "× なし", "○ JCS: 予約照会APIあり"],
    ["法人割引・法人料金", "× 一般料金のみ", "○ 法人一括決済プランの料金表示"],
    ["管理の統一性", "○ 全サービス同一運用", "△ 宿泊だけ別運用（管理画面あり）"],
    ["費用", "○ 完全無料", "○ JCS/Raccoとも無料"],
], Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.2))

# Existing system note
note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(6.2), Inches(1.2))
note_box.fill.solid()
note_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xEE)
note_box.line.color.rgb = GREEN
note_box.line.width = Pt(2)
add_text(slide, "既存の仕組みでPJ紐付けは解決済み", Inches(0.8), Inches(6.0), Inches(5.6), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "Slack /trip でHubSpot案件番号を入力（申請時）",
    "MFカード明細とカード予測レコードで自動突合",
    "→ 法人契約なしでもPJ別コスト管理が可能",
], Inches(0.8), Inches(6.35), Inches(5.6), Inches(0.7), size=12, color=DARK_TEXT)

# When to upgrade box
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(5.9), Inches(5.8), Inches(1.2))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
box2.line.color.rgb = ORANGE
box2.line.width = Pt(2)
add_text(slide, "案Bへの移行タイミング", Inches(7.3), Inches(6.0), Inches(5.2), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "宿泊費の上限管理をしたくなった時",
    "宿名・日程を含む詳細CSVが必要になった時",
    "API連携で自動化したくなった時",
], Inches(7.3), Inches(6.35), Inches(5.2), Inches(0.7), size=12, color=DARK_TEXT)

# =====================================================
# Slide 5: JCS vs Racco (参考: 案Bの場合)
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "参考: 案Bで法人契約する場合 — JCS vs Racco", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=26, color=WHITE, bold=True)

styled_table(slide, [
    ["比較項目", "じゃらんJCS", "楽天Racco", "評価"],
    ["審査", "リクルート所定の審査あり（期間未確認）", "与信審査2週間+ANA Biz 1.5ヶ月", "要確認"],
    ["API連携", "予約照会APIあり", "おそらく不可", "JCS優位"],
    ["退職時対応", "個人ID（影響小）", "法人PW変更が必要", "JCS優位"],
    ["宿泊施設数", "じゃらんの豊富な在庫", "楽天トラベルの在庫", "同等"],
    ["航空券予約", "不可（別手配）", "パッケージ商品で可", "Racco優位"],
    ["請求書", "紙郵送のみ", "1枚にまとまる", "Racco優位"],
    ["管理ルール", "上限金額・金券除外設定可", "リアルタイム利用把握", "同等"],
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.0))

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
box.line.color.rgb = ORANGE
box.line.width = Pt(2)

add_text(slide, "案Bに移行する場合はJCSが優位（API・退職対応）。Raccoの航空券優位は個人アカウント運用で不要",
         Inches(1.2), Inches(6.0), Inches(11), Inches(0.8), size=16, color=PRIMARY, bold=True)

# =====================================================
# Slide 5: Operation Flow
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "運用フロー", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

steps = [
    ("① 事前承認", ["Slack /trip", "コマンドで申請", "上長がSlack承認"], ACCENT),
    ("② 宿泊手配", ["じゃらんJCSで予約", "法人一括決済", "PJ番号を入力"], GREEN),
    ("③ 航空券手配", ["ANA Biz or", "JALオンライン", "MFカード決済"], GREEN),
    ("④ 出張実施", ["新幹線: スマートEX", "レンタカー: 各社", "MFカード決済"], GRAY_TEXT),
    ("⑤ 月次精算", ["JCS: 確定CSV取込", "航空券: MFカード自動取込", "新幹線等: MFカード自動取込"], ORANGE),
]

cw = Inches(2.15)
ch = Inches(2.2)
for i, (title, lines, color) in enumerate(steps):
    x = Inches(0.6) + i * (cw + Inches(0.2))
    y = Inches(1.6)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = TABLE_BORDER
    card.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, title, x + Inches(0.15), y + Inches(0.2), cw - Inches(0.3), Inches(0.4),
             size=16, color=PRIMARY, bold=True)
    add_bullets(slide, lines, x + Inches(0.15), y + Inches(0.65), cw - Inches(0.3), Inches(1.4),
                size=13, color=DARK_TEXT)
    if i < len(steps) - 1:
        add_text(slide, "→", x + cw + Inches(0.02), y + ch / 2 - Inches(0.15), Inches(0.2), Inches(0.3),
                 size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Monthly settlement detail
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.8))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.fill.background()

add_text(slide, "月次精算の詳細（管理本部）", Inches(1.0), Inches(4.4), Inches(10), Inches(0.4),
         size=18, color=PRIMARY, bold=True)

styled_table(slide, [
    ["手配先", "データ取得", "MF登録方法", "証憑"],
    ["JCS（宿泊）", "翌月8日 確定CSV", "CSV加工 → MF経費一括登録", "紙請求書スキャン → MF添付"],
    ["ANA Biz/JAL", "MFカード明細自動取込", "自動（MF連携済み）", "利用明細"],
    ["スマートEX", "MFカード明細自動取込", "自動（MF連携済み）", "利用明細"],
    ["レンタカー", "MFカード明細自動取込", "自動（MF連携済み）", "利用明細"],
], Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.0))

# =====================================================
# Slide: 3つの購入フロー（通常・概算・緊急事後）
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "購入フロー: 通常 / 概算 / 緊急事後", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=28, color=WHITE, bold=True)

# 3 flow cards
fw = Inches(3.8)
fh = Inches(3.6)
fy = Inches(1.3)

add_phase_card(slide, Inches(0.5), fy, fw, fh, "通常フロー", [
    "1. /purchase or /trip で事前申請",
    "2. 上長がSlackで承認",
    "3. 各サイトでMFカード購入",
    "4. MFカード明細で自動取込",
    "5. カード予測で申請と自動突合",
], GREEN)

add_phase_card(slide, Inches(4.75), fy, fw, fh, "概算フロー（金額未確定）", [
    "1. /purchase で概算額を入力",
    "   「概算」フラグ付きで申請",
    "2. 上長が概算ベースで承認",
    "3. MFカードで購入（実額）",
    "4. 実額と概算を自動比較",
    "5. 差額大なら再承認通知",
], ACCENT)

add_phase_card(slide, Inches(9.0), fy, fw, fh, "緊急事後フロー", [
    "1. MFカードで緊急購入",
    "   （事前申請する暇がない）",
    "2. 24h以内に事後報告",
    "   品目・理由・緊急理由を入力",
    "3. 上長が事後承認",
    "4. 未報告は週次でアラート",
], ORANGE)

# Bottom comparison: MF card vs cash
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.fill.background()

add_text(slide, "なぜ「緊急時こそMFカード」か — 現金立替との比較", Inches(0.8), Inches(5.3), Inches(10), Inches(0.35),
         size=16, color=PRIMARY, bold=True)

styled_table(slide, [
    ["観点", "現金立替", "MFカード（物理）"],
    ["購入の記録", "本人の申告のみ", "カード明細に自動記録"],
    ["未申告の検知", "不可能", "明細に申請紐付けなし → アラート"],
    ["金額の正確性", "領収書頼み", "明細が正（改ざん不可）"],
    ["社員の立替負担", "あり", "なし"],
    ["実店舗での利用", "○", "○ 物理カードで対応"],
], Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.4))

# =====================================================
# Slide 6: System Integration Roadmap
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "システム連携ロードマップ", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

pw = Inches(3.7)
ph = Inches(3.8)
add_phase_card(slide, Inches(0.6), Inches(1.5), pw, ph, "Phase 1: 手動運用（5月~）", [
    "JCS CSV → 手動加工 → MF経費一括登録",
    "ANA Biz/JAL → MFカード明細で自動取込",
    "MFカード → 自動取込（新幹線・レンタカー）",
    "購買管理システムの仕訳管理画面で一括登録",
], GREEN)
add_phase_card(slide, Inches(4.55), Inches(1.5), pw, ph, "Phase 2: API連携（6月~）", [
    "JCS予約照会API → 購買管理で自動取込",
    "Slack通知（予約完了・変更・キャンセル）",
    "MF経費への自動登録",
    "PJ番号による部門別集計の自動化",
], ACCENT)
add_phase_card(slide, Inches(8.5), Inches(1.5), pw, ph, "Phase 3: 統合管理（将来構想）", [
    "全出張データを購買管理システムに集約",
    "出張費の予実管理ダッシュボード",
    "部門別・PJ別の出張コスト可視化",
    "承認ワークフローの統合",
], ORANGE)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
box.line.color.rgb = ORANGE
box.line.width = Pt(2)

add_text(slide, "JCSの予約照会APIは Phase 2 以降の自動化において重要な差別化要素。Raccoにはこの選択肢がない。",
         Inches(1.0), Inches(5.9), Inches(11), Inches(0.8), size=16, color=PRIMARY, bold=True)

# =====================================================
# Slide 7: EC/購買サイト MFバーチャルカード対応一覧
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "EC/購買サイト × MFバーチャルカード 対応一覧", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=26, color=WHITE, bold=True)

styled_table(slide, [
    ["サイト", "カード決済(VISA)", "法人アカウント", "個人ユーザー", "承認フロー", "備考"],
    ["Amazon Business", "○", "○", "○ 配下に作成", "○ 最大6段階", "Business Prime ¥5,900~/年"],
    ["モノタロウ", "○", "○", "○ エンタープライズ版", "○ 多段階", "エンタープライズは要問合せ"],
    ["アスクル", "○", "○", "○", "△ 限定的", "無料"],
    ["たのめーる", "○", "○", "○ 営業経由", "○", "大塚商会と要相談"],
    ["カウネット", "○ 3Dセキュア必須", "○", "○", "△ 承認時カード不可", "無料"],
    ["ビックカメラ法人", "○ マルチペイメント", "○", "△", "△", "無料"],
    ["ヨドバシ法人", "× 振込のみ", "△", "×", "×", "個人アカウントならカード可"],
    ["いだてん", "要確認", "○", "—", "—", "プラス社に確認"],
], Inches(0.4), Inches(1.3), Inches(12.5), Inches(3.8))

# Summary boxes
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.4), Inches(6.0), Inches(1.7))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xEE)
box1.line.color.rgb = GREEN
box1.line.width = Pt(2)
add_text(slide, "MFバーチャルカードの強み", Inches(0.7), Inches(5.5), Inches(5.4), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "社員ごとに個別番号を即時発行（VISA）",
    "カードごとに利用上限を管理者が設定可能",
    "用途別・PJ別にもカード発行可能",
    "→ ECサイト側の購買制御がなくてもカード側で制御",
], Inches(0.7), Inches(5.85), Inches(5.4), Inches(1.1), size=12, color=DARK_TEXT)

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(5.4), Inches(6.2), Inches(1.7))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
box2.line.color.rgb = ORANGE
box2.line.width = Pt(2)
add_text(slide, "注意点・要確認", Inches(7.0), Inches(5.5), Inches(5.6), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "ヨドバシ法人: カード不可 → 個人アカウント運用で対応",
    "いだてん: カード決済対応をプラス社に確認",
    "カウネット: 承認フロー利用時はカード決済不可",
    "Amazon Business: 高度機能はBusiness Prime契約が必要",
], Inches(7.0), Inches(5.85), Inches(5.6), Inches(1.1), size=12, color=DARK_TEXT)

# =====================================================
# Slide: 適格請求書（インボイス）・証憑管理
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "適格請求書（インボイス）・証憑管理", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=28, color=WHITE, bold=True)

# Key message
alert = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.8))
alert.fill.solid()
alert.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
alert.line.color.rgb = RED
alert.line.width = Pt(2)
add_text(slide, "MFカード明細は適格請求書ではない — 各購入先からの適格請求書（領収書）の取得・保存は引き続き必須",
         Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5), size=15, color=RED, bold=True)

# Table: invoice acquisition per service
styled_table(slide, [
    ["購入先", "適格請求書の形態", "取得方法", "保存"],
    ["Amazon", "電子領収書（PDF）", "注文履歴からDL", "電子保存"],
    ["じゃらん/楽天トラベル", "宿泊明細書", "マイページからDL", "電子保存"],
    ["ANA / JAL", "搭乗明細・eチケット控え", "マイページからDL", "電子保存"],
    ["スマートEX / えきねっと", "利用明細・領収書", "Web/アプリから取得", "電子保存"],
    ["モノタロウ / アスクル", "電子請求書", "マイページからDL", "電子保存"],
    ["実店舗（緊急購入）", "紙レシート", "社員が保管・提出", "スキャン → 電子保存"],
], Inches(0.5), Inches(2.2), Inches(7.0), Inches(2.8))

# Right side: accounting treatment
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.2), Inches(5.0), Inches(2.8))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.fill.background()

add_text(slide, "仕訳における取引先の扱い", Inches(8.0), Inches(2.3), Inches(4.6), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "カード払いでも取引先コードを設定",
    "  → 適格請求書の発行元を管理",
    "  → 購入先別の支出分析に必要",
    "",
    "仕訳パターン:",
    "  借方: 費用科目（消耗品費等）",
    "  貸方: 未払金（MFカード:未請求）",
    "        取引先: 購入先名（Amazon等）",
], Inches(8.0), Inches(2.65), Inches(4.6), Inches(2.2), size=12, color=DARK_TEXT)

# Bottom boxes
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(6.0), Inches(1.8))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
box1.line.color.rgb = ACCENT
box1.line.width = Pt(2)
add_text(slide, "運用ルール", Inches(0.8), Inches(5.4), Inches(5.4), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "購入後、適格請求書をSlackスレッドに添付（既存フロー）",
    "電子保存対応: PDF/スクリーンショットをスレッドに投稿",
    "実店舗レシート: スマホ撮影 → スレッドに投稿",
    "証憑未提出は週次で管理本部がフォロー（既存）",
], Inches(0.8), Inches(5.75), Inches(5.4), Inches(1.2), size=12, color=DARK_TEXT)

box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.3), Inches(6.0), Inches(1.8))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
box2.line.color.rgb = ORANGE
box2.line.width = Pt(2)
add_text(slide, "免税事業者リスク（インボイス経過措置）", Inches(7.1), Inches(5.4), Inches(5.4), Inches(0.35),
         size=15, color=PRIMARY, bold=True)
add_bullets(slide, [
    "インボイス未登録事業者からの購入:",
    "  ~2026/9: 仕入税額の80%が控除可能",
    "  2026/10~2029/9: 50%控除",
    "  2029/10~: 控除不可",
    "→ 購入先の登録番号(T+13桁)を確認する運用を推奨",
], Inches(7.1), Inches(5.75), Inches(5.4), Inches(1.2), size=12, color=DARK_TEXT)

# =====================================================
# Slide 8: Schedule (案Aベース)
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "導入スケジュール", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

styled_table(slide, [
    ["時期", "アクション", "担当", "備考"],
    ["4月前半", "MFカード全社員配布（物理+バーチャル）", "管理本部", "未配布社員への追加発行。物理=実店舗・緊急用"],
    ["4月前半", "各サービスの会社用個人アカウント作成", "各社員", "出張: じゃらん/ANA/JAL/スマートEX/えきねっと"],
    ["4月前半", "Amazon Business法人アカウント設定", "管理本部", "個人ユーザー追加・購買ポリシー設定"],
    ["4月中旬", "社員向け利用ガイド作成・周知", "管理本部", "Slack /trip の案内・各サイトの使い方"],
    ["5月~", "運用開始（案A: 個人アカウント統一）", "全社", "法人契約なしでスタート"],
    ["随時", "法人契約の追加検討（案B移行判断）", "管理本部", "JCS/エクスプレス予約等、必要に応じて"],
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.2))

add_text(slide, "※ 法人契約なしのため審査不要。アカウント作成+MFカード登録で即日利用開始可能",
         Inches(0.8), Inches(5.0), Inches(11), Inches(0.4), size=14, color=GRAY_TEXT)

# =====================================================
# Slide 9: Risks & Cost (案Aベース)
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "リスクと対策 / コスト", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

styled_table(slide, [
    ["リスク", "影響度", "対策"],
    ["宿泊費の詳細把握が困難（宿名等なし）", "中", "Slack /trip で行先・宿泊先を申告。必要時に案Bへ移行"],
    ["社員の予約に上限がない", "低", "MFバーチャルカードの利用上限で制御可能"],
    ["ヨドバシ法人がカード決済不可", "低", "個人アカウント+MFカードで対応、または振込"],
    ["いだてんのカード決済対応が不明", "低", "プラス社に確認。不可なら請求書払い"],
    ["MF会計Plus OAuth認証が未実施", "高", "API連携コードは準備済み。OAuth初回認証を早期に実施"],
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.8))

add_text(slide, "コスト", Inches(0.8), Inches(4.6), Inches(10), Inches(0.4),
         size=22, color=PRIMARY, bold=True)

styled_table(slide, [
    ["項目", "金額", "備考"],
    ["出張手配（全サービス）", "無料", "個人アカウント+MFカード、法人契約なし"],
    ["Amazon Business Prime", "¥5,900~/年", "高度な購買管理が必要な場合のみ"],
    ["MFバーチャルカード", "既存契約内", "追加費用なし（MFビジネスカード契約済み前提）"],
    ["法人契約追加時（JCS等）", "無料", "将来の案B移行時も追加費用なし"],
], Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.0))

# =====================================================
# Slide 10: Next Actions
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide)
add_text(slide, "次のアクション", Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
         size=30, color=WHITE, bold=True)

actions = [
    ("即時", [
        "本提案書の社内レビュー・承認",
        "MFカード全社員配布（物理+バーチャル）",
        "Amazon Business法人アカウント設定",
    ], GREEN),
    ("4月中", [
        "各サービスの会社用アカウント作成",
        "いだてんカード決済対応をプラス社に確認",
        "社員向け利用ガイド作成・周知",
    ], ACCENT),
    ("5月以降", [
        "案A運用開始（個人アカウント統一）",
        "運用状況を見て案B移行を検討",
        "MF会計Plus OAuth初回認証の実施",
    ], ORANGE),
]

cw = Inches(3.7)
ch = Inches(4.0)
for i, (timing, items, color) in enumerate(actions):
    x = Inches(0.6) + i * (cw + Inches(0.25))
    y = Inches(1.5)
    add_phase_card(slide, x, y, cw, ch, timing, [f"☐ {item}" for item in items], color)

add_text(slide, "ご検討よろしくお願いいたします", Inches(0.8), Inches(6.2), Inches(12), Inches(0.5),
         size=18, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# Save
output = "docs/travel-services/出張手配サービス導入_推奨案.pptx"
prs.save(output)
print(f"Saved: {output} ({len(prs.slides)} slides)")

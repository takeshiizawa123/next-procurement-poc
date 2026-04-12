"""業務フロー設計書 — 統合経費管理（購買・出張・立替）"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0F, 0x0F, 0x1A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK = RGBColor(0xEC, 0x48, 0x99)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SURFACE = RGBColor(0x1A, 0x1C, 0x2E)


def dbg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG


def hero(title, sub="", accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dbg(s)
    al = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.6), Inches(1.2), Pt(4))
    al.fill.solid()
    al.fill.fore_color.rgb = accent
    al.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(2.5))
    t.text_frame.word_wrap = True
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(48)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    if sub:
        p2 = t.text_frame.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(20)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(20)


def sec(num, title, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dbg(s)
    nt = s.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(3), Inches(3))
    nt.text_frame.paragraphs[0].text = f"{num:02d}"
    nt.text_frame.paragraphs[0].font.size = Pt(120)
    nt.text_frame.paragraphs[0].font.bold = True
    nt.text_frame.paragraphs[0].font.color.rgb = accent
    tt = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1.5))
    tt.text_frame.paragraphs[0].text = title
    tt.text_frame.paragraphs[0].font.size = Pt(40)
    tt.text_frame.paragraphs[0].font.bold = True
    tt.text_frame.paragraphs[0].font.color.rgb = WHITE


def title_slide(title, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dbg(s)
    t = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(28)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    al = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(0.8), Pt(3))
    al.fill.solid()
    al.fill.fore_color.rgb = accent
    al.line.fill.background()
    return s


def bullets(title, items, accent=CYAN):
    s = title_slide(title, accent)
    ct = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(12), Inches(5.5))
    ct.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = ct.text_frame.paragraphs[0] if i == 0 else ct.text_frame.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.level = item[1]
            p.font.size = Pt(16 if item[1] > 0 else 22)
            p.font.color.rgb = MUTED if item[1] > 0 else LIGHT
        elif item == "":
            p.text = ""
            p.font.size = Pt(8)
        else:
            p.text = item
            p.font.size = Pt(22)
            p.font.color.rgb = LIGHT
        p.space_after = Pt(6)


def box(s, x, y, w, h, fill, text="", text_color=WHITE, font_size=14, bold=False, border=None):
    """ボックス描画ヘルパー"""
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bx.fill.solid()
    bx.fill.fore_color.rgb = fill
    if border:
        bx.line.color.rgb = border
        bx.line.width = Pt(1.5)
    else:
        bx.line.fill.background()
    if text:
        bx.text_frame.word_wrap = True
        bx.text_frame.margin_left = Inches(0.1)
        bx.text_frame.margin_right = Inches(0.1)
        bx.text_frame.margin_top = Inches(0.05)
        bx.text_frame.margin_bottom = Inches(0.05)
        bx.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = bx.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
    return bx


def label(s, x, y, w, h, text, font_size=12, color=LIGHT, bold=False, align=PP_ALIGN.LEFT):
    """ラベル描画ヘルパー"""
    t = s.shapes.add_textbox(x, y, w, h)
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return t


def arrow(s, x1, y1, x2, y2, color=CYAN):
    """矢印描画ヘルパー"""
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.5)
    line.line.end_arrow_type = 2  # arrow head
    return line


def tbl(title, headers, rows, accent=CYAN, top=Inches(1.7)):
    s = title_slide(title, accent)
    cols = len(headers)
    n = len(rows) + 1
    ts = s.shapes.add_table(n, cols, Inches(0.5), top, Inches(12.3), Inches(0.45) * n)
    table = ts.table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = SURFACE
        for par in c.text_frame.paragraphs:
            par.font.size = Pt(14)
            par.font.bold = True
            par.font.color.rgb = accent
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = table.cell(i + 1, j)
            c.text = str(v)
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else BG
            for par in c.text_frame.paragraphs:
                par.font.size = Pt(13)
                par.font.color.rgb = LIGHT


# ==================== SLIDES ====================

# Part 0: タイトル

hero("統合経費管理 業務フロー設計", "購買・出張・立替を一元管理\nMF経費 + MF会計Plus 連携設計（B案）")

# 用語整理
tbl("システムの主要コンポーネント", ["コンポーネント", "役割"], [
    ["購買管理システム", "申請・承認・照合・仕訳作成のすべての司令塔（自社開発・Vercel）"],
    ["MFビジネスカード", "従業員1人1枚のバーチャルカード。すべての経費決済に使用"],
    ["MF経費", "カード明細の取り込み口。購買管理がAPIで明細を取得（ユーザーは触らない）"],
    ["MF会計Plus", "仕訳の最終出力先。購買管理がAPIで仕訳を直接登録"],
    ["Slack", "通知・承認の窓口（DMボタンで承認、申請完了通知等）"],
    ["Webアプリ", "申請・マイページ・管理画面のメインUI"],
])

# Part 1: 全体像

sec(1, "全体像 — システム構成", CYAN)

# 全体構成図
s = title_slide("全体構成 — データの流れ", CYAN)

# Layer 1: User
label(s, Inches(0.5), Inches(1.7), Inches(2), Inches(0.4), "Layer 1: ユーザー操作", font_size=14, color=CYAN, bold=True)
box(s, Inches(0.5), Inches(2.1), Inches(3.5), Inches(0.7), PURPLE, "Webアプリ\n申請・マイページ・管理画面", font_size=12, bold=True)
box(s, Inches(4.2), Inches(2.1), Inches(3.5), Inches(0.7), PINK, "Slack\n承認DMボタン・通知", font_size=12, bold=True)

# Layer 2: 購買管理システム（中央）
label(s, Inches(0.5), Inches(3.0), Inches(8), Inches(0.4), "Layer 2: 購買管理システム（司令塔）", font_size=14, color=CYAN, bold=True)
box(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.2), CARD, "申請・承認・予測テーブル・照合エンジン・仕訳作成", font_size=14, bold=True, border=CYAN)

# Layer 3: 外部システム
label(s, Inches(0.5), Inches(4.8), Inches(6), Inches(0.4), "Layer 3: 外部システム（バックエンド連携のみ）", font_size=14, color=CYAN, bold=True)
box(s, Inches(0.5), Inches(5.2), Inches(3.8), Inches(1.5), GREEN, "MFビジネスカード\n各従業員に1枚\nすべての決済に使用", font_size=12, bold=True)
box(s, Inches(4.6), Inches(5.2), Inches(3.8), Inches(1.5), ORANGE, "MF経費\nカード明細の取り込み口\nAPI で明細取得", font_size=12, bold=True)
box(s, Inches(8.7), Inches(5.2), Inches(4.0), Inches(1.5), CYAN, "MF会計Plus\n仕訳の最終出力先\nAPI で仕訳登録", font_size=12, bold=True)

# 矢印（簡易）
label(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
      "決済 → MF経費（自動取込）→ 購買管理がAPI取得 → 照合 → MF会計Plus に仕訳登録",
      font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Part 2: 3つの経費パターン
sec(2, "3つの経費パターン", PURPLE)

s = title_slide("経費の3パターン — すべて統一フロー", PURPLE)
patterns = [
    ("購買", "事前申請して物品・サービスを購入", "Amazon, モノタロウ等\nバーチャルカード決済", PURPLE),
    ("出張", "事前申請して出張に行く", "スマートEX, じゃらん, ANA等\nバーチャルカード決済", PINK),
    ("立替", "やむを得ず個人で立替→精算", "現金払い, 個人カード\n後日会社が振込", ORANGE),
]
for i, (label_text, desc, ex, clr) in enumerate(patterns):
    x = Inches(0.5) + i * Inches(4.3)
    box(s, x, Inches(2.0), Inches(4.0), Inches(0.7), clr, label_text, font_size=22, bold=True)
    box(s, x, Inches(2.85), Inches(4.0), Inches(1.5), CARD, desc, font_size=14, text_color=LIGHT, border=clr)
    box(s, x, Inches(4.5), Inches(4.0), Inches(1.5), SURFACE, ex, font_size=12, text_color=MUTED)

label(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
      "▶ いずれもWebアプリ/Slackで申請 → 部門長承認 → カード決済 → 自動照合 → 仕訳",
      font_size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

# Part 3: 申請フロー
sec(3, "申請フロー", PINK)

s = title_slide("申請フロー — Webアプリ & Slack", PINK)

# 縦並びの5ステップ
steps_data = [
    ("①", "申請", "Webアプリ or Slackで申請\n(購買/出張/立替を選択)", PURPLE),
    ("②", "予測テーブル登録", "申請内容をDBに保存\n(従業員ID + 金額 + 日付)", CYAN),
    ("③", "部門長承認", "Slack DMで [承認] [差戻し] ボタン\n10秒で完了", PINK),
    ("④", "申請者通知", "承認完了をSlackで通知\n申請者は購入実行へ", GREEN),
    ("⑤", "予測テーブル状態更新", "「approved」に変更\nカード決済を待つ状態", ORANGE),
]
y = Inches(1.8)
for i, (num, title, desc, clr) in enumerate(steps_data):
    x = Inches(0.5)
    box(s, x, y, Inches(0.9), Inches(0.8), clr, num, font_size=24, bold=True)
    box(s, Inches(1.5), y, Inches(3.0), Inches(0.8), CARD, title, font_size=16, bold=True, border=clr)
    box(s, Inches(4.6), y, Inches(8.2), Inches(0.8), SURFACE, desc, font_size=13, text_color=LIGHT)
    y += Inches(1.0)

# Part 4: 決済とカード明細取得
sec(4, "カード決済 → 明細取得", GREEN)

s = title_slide("決済からMF経費取込まで", GREEN)

# 縦並びフロー
flow = [
    ("従業員", "バーチャルカードで決済\n(Amazon/スマートEX/じゃらん等)", CYAN),
    ("MFビジネスカード", "決済データ自動連携", GREEN),
    ("MF経費", "カード明細として取り込み\noffice_member_id で従業員を自動紐付け", ORANGE),
    ("購買管理システム", "fetchCardStatements API で取得\n(週次cron + 任意手動実行)", CYAN),
]
y = Inches(1.8)
for i, (who, what, clr) in enumerate(flow):
    box(s, Inches(0.5), y, Inches(3.5), Inches(1.0), clr, who, font_size=18, bold=True)
    box(s, Inches(4.2), y, Inches(8.6), Inches(1.0), CARD, what, font_size=14, border=clr)
    if i < len(flow) - 1:
        # 下向き矢印
        arrow_lbl = s.shapes.add_textbox(Inches(2.0), y + Inches(1.0), Inches(0.5), Inches(0.3))
        arrow_lbl.text_frame.paragraphs[0].text = "▼"
        arrow_lbl.text_frame.paragraphs[0].font.size = Pt(20)
        arrow_lbl.text_frame.paragraphs[0].font.color.rgb = MUTED
    y += Inches(1.4)

# Part 5: 照合フロー
sec(5, "照合フロー", ORANGE)

s = title_slide("照合エンジン — 予測テーブル × カード明細", ORANGE)

# 左右レイアウト
# 左: 予測テーブル
box(s, Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.6), PURPLE, "予測テーブル（申請データ）", font_size=14, bold=True)
preds = [
    "PO-202604-0001 / 田中 / ¥14,000 / 2026-04-15 / スマートEX (出張)",
    "PO-202604-0002 / 田中 / ¥30,000 / 2026-04-15 / じゃらん (出張)",
    "PO-202604-0003 / 鈴木 / ¥3,500 / 2026-04-15 / Amazon (購買)",
    "PO-202604-0004 / 山田 / ¥8,000 / 2026-04-16 / モノタロウ (購買)",
]
for i, p in enumerate(preds):
    box(s, Inches(0.5), Inches(2.5) + Inches(0.6) * i, Inches(5.5), Inches(0.55),
        SURFACE, p, font_size=10, text_color=LIGHT)

# 右: カード明細
box(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.6), ORANGE, "MF経費カード明細", font_size=14, bold=True)
stmts = [
    "EX東海 / ¥14,000 / 2026-04-15 / member:田中",
    "じゃらんnet / ¥29,800 / 2026-04-15 / member:田中",
    "Amazon.co.jp / ¥3,500 / 2026-04-15 / member:鈴木",
    "MonotaRO / ¥8,000 / 2026-04-16 / member:山田",
]
for i, st in enumerate(stmts):
    box(s, Inches(7.0), Inches(2.5) + Inches(0.6) * i, Inches(5.8), Inches(0.55),
        SURFACE, st, font_size=10, text_color=LIGHT)

# 中央の照合キー
box(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.7), CARD,
    "照合キー: office_member_id (従業員) × 金額 × 日付", font_size=14, bold=True, border=ORANGE)

# 照合結果
label(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4), "照合結果", font_size=14, color=GREEN, bold=True)
box(s, Inches(0.5), Inches(6.5), Inches(3.0), Inches(0.7), GREEN, "完全一致\n→ 自動仕訳", font_size=11, bold=True)
box(s, Inches(3.7), Inches(6.5), Inches(3.0), Inches(0.7), ORANGE, "差額あり\n→ 調整仕訳", font_size=11, bold=True)
box(s, Inches(6.9), Inches(6.5), Inches(3.0), Inches(0.7), PURPLE, "申請なし\n→ 未申請利用アラート", font_size=11, bold=True)
box(s, Inches(10.1), Inches(6.5), Inches(2.7), Inches(0.7), RED, "明細なし\n→ 督促", font_size=11, bold=True)

# Part 6: 仕訳
sec(6, "仕訳の作成", CYAN)

s = title_slide("仕訳作成 — 購買管理システムが直接生成", CYAN)

# 説明
label(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
      "MF会計Plusの自動仕訳ルールは使わず、購買管理が直接APIで仕訳を作成（B案）",
      font_size=14, color=MUTED)

# 例: 出張の仕訳
label(s, Inches(0.5), Inches(2.4), Inches(12), Inches(0.4), "例: 出張交通費の仕訳", font_size=16, color=CYAN, bold=True)

# Stage 1
box(s, Inches(0.5), Inches(2.9), Inches(2.0), Inches(0.5), PURPLE, "Stage 1", font_size=14, bold=True)
box(s, Inches(2.7), Inches(2.9), Inches(10.1), Inches(0.5), CARD,
    "借: 旅費交通費 14,000 / 貸: 未払金:MFカード:未請求 14,000", font_size=12)
label(s, Inches(2.7), Inches(3.4), Inches(10), Inches(0.3),
      "▶ 出張申請承認時に作成（経費認識）", font_size=11, color=MUTED)

# Stage 2
box(s, Inches(0.5), Inches(3.9), Inches(2.0), Inches(0.5), CYAN, "Stage 2", font_size=14, bold=True)
box(s, Inches(2.7), Inches(3.9), Inches(10.1), Inches(0.5), CARD,
    "借: 未払金:MFカード:未請求 14,000 / 貸: 未払金:MFカード:請求 14,000", font_size=12)
label(s, Inches(2.7), Inches(4.4), Inches(10), Inches(0.3),
      "▶ カード明細とマッチング後に作成（請求確定）", font_size=11, color=MUTED)

# Stage 3
box(s, Inches(0.5), Inches(4.9), Inches(2.0), Inches(0.5), GREEN, "Stage 3", font_size=14, bold=True)
box(s, Inches(2.7), Inches(4.9), Inches(10.1), Inches(0.5), CARD,
    "借: 未払金:MFカード:請求 14,000 / 貸: 普通預金 14,000", font_size=12)
label(s, Inches(2.7), Inches(5.4), Inches(10), Inches(0.3),
      "▶ 銀行引落時に作成（支払完了）", font_size=11, color=MUTED)

# 摘要付加の説明
box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2), SURFACE,
    "摘要は申請データと結合: 「東京→大阪 新幹線のぞみ ¥14,000 申請者:田中 (PO-202604-0001)」\nMF会計Plus上で経理が確認しやすい形式に",
    font_size=12, text_color=LIGHT, border=CYAN)

# Part 7: B案の特徴
sec(7, "B案の特徴と運用ポイント", PINK)

bullets("B案 (購買管理システムがStage 2を作る) のメリット", [
    "自動仕訳ルールの設定が不要 — 経理担当者の負担ゼロ",
    "新規加盟店が出ても自動対応 — 運用保守がほぼ不要",
    "仕訳の摘要・補助科目を自由にカスタマイズ可能",
    "  ▶ 出張申請の詳細（区間・宿泊先等）を摘要に付加",
    "  ▶ 部門・PJ も申請データから自動設定",
    "二重仕訳のリスクなし — 仕訳作成は購買管理だけ",
    "MF経費 → 購買管理 → MF会計Plus の経路が一本化",
    "office_member_id で従業員を確実に特定（card_last4不要）",
])

bullets("B案で必要な事前設定", [
    "MFビジネスカード × MF経費 の自動連携設定",
    "  ▶ 各従業員のバーチャルカードをMF経費アカウントに紐付け",
    "  ▶ 1人1枚のため設定はシンプル",
    "MF会計Plusの自動仕訳ルールは「使わない（OFF設定）」",
    "  ▶ 二重仕訳防止のため必須",
    "MF経費APIアクセストークンの発行",
    "MF会計Plus OAuth認証",
    "従業員マスタとMF経費 office_member の紐付け確認",
])

# Part 8: 例外処理
sec(8, "例外フロー", RED)

tbl("想定される例外と対応", ["パターン", "発生例", "対応方針"], [
    ["申請なしの利用", "出張中に急遽タクシー利用", "管理者にアラート → 事後申請を促す"],
    ["概算と実額の差", "新幹線予約変更で金額変更", "差額±5%以内は自動承認、超える場合は管理者確認"],
    ["申請したのに明細なし", "予約後にキャンセル", "1週間後に申請者に確認DM"],
    ["重複決済", "同日同額の購買と出張", "office_member_id + 申請内容で個別マッチング"],
    ["返金・キャンセル", "出張キャンセルでカード返金", "負の値の明細として取り込み、Stage 1の取消仕訳を作成"],
    ["不明な明細", "従業員特定不可", "管理者ダッシュボードに「未照合」として表示、手動対応"],
])

# Part 9: 役割分担
sec(9, "役割分担", GREEN)

tbl("関係者ごとの作業", ["役割", "やること", "やらないこと"], [
    ["申請者(従業員)", "Webアプリ/Slackで申請 → カード決済 → 検収", "MF経費・MF会計Plusに触る"],
    ["承認者(部門長)", "Slack DMで承認/差戻し", "個別の照合・仕訳確認"],
    ["管理本部", "未照合明細の確認、例外対応、月次締め", "通常の仕訳作成・自動仕訳ルール管理"],
    ["購買管理システム", "申請受付・予測登録・照合・仕訳自動生成", "—"],
    ["MFビジネスカード", "決済・明細をMF経費に連携", "—"],
    ["MF経費", "カード明細の保管・API提供", "申請承認・仕訳作成"],
    ["MF会計Plus", "仕訳の最終保管・経理出力", "自動仕訳ルール（OFF）"],
])

# Part 10: 今後の検証ポイント
sec(10, "今後の検証ポイント", ORANGE)

bullets("実装前に確認すべきこと", [
    "✓ MF経費APIで office_member_id が取得可能（実機確認済み）",
    "□ MFビジネスカード → MF経費 の自動連携設定方法",
    "  ▶ 1人1枚の紐付けが正しく動くか",
    "□ カード明細の反映タイミング（即時 or 翌営業日）",
    "□ 同日同額の複数決済の区別方法",
    "□ 返金・キャンセル時の明細形式",
    "□ MF会計Plus の自動仕訳ルールをOFFにできるか",
    "□ 二重計上防止のための一意キー設計",
    "□ 立替精算のフロー（MF経費を通すか直接仕訳か）",
])

# Part 11: 段階的導入計画
sec(11, "段階的導入計画", PURPLE)

s = title_slide("Phase別ロードマップ", PURPLE)

phases = [
    ("Phase 1", "基盤構築", "従業員マスタ拡張・予測テーブル統一・出張申請Web化", PURPLE),
    ("Phase 2", "MF経費連携", "fetchCardStatements改修・office_member_id ベース照合", CYAN),
    ("Phase 3", "出張承認", "/trip 承認フロー追加・部門長DMボタン", PINK),
    ("Phase 4", "立替Web化", "/expense/new 新規・OCR連携", GREEN),
    ("Phase 5", "Stage 2自動化", "購買管理がMF会計Plus APIでStage 2仕訳作成", ORANGE),
    ("Phase 6", "本番展開", "管理本部→他部門の段階展開・運用マニュアル整備", RED),
]
y = Inches(1.8)
for i, (phase, name, desc, clr) in enumerate(phases):
    box(s, Inches(0.5), y, Inches(1.8), Inches(0.7), clr, phase, font_size=14, bold=True)
    box(s, Inches(2.5), y, Inches(2.5), Inches(0.7), CARD, name, font_size=14, bold=True, border=clr)
    box(s, Inches(5.2), y, Inches(7.6), Inches(0.7), SURFACE, desc, font_size=12, text_color=LIGHT)
    y += Inches(0.85)

# Part 12: クロージング
hero("ご検討よろしくお願いします", "本資料は B案（購買管理システム主導）の業務フロー設計\n実装前に運用詳細・例外パターンをすり合わせましょう", PURPLE)

# 保存
out_path = os.path.join(os.path.dirname(__file__), "..", "workflow-design-b-route.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")

"""購買管理・出張管理 想定運用ガイド — Silicon Valley Style PPT (Final)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0F, 0x0F, 0x1A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
SURFACE = RGBColor(0x1A, 0x1C, 0x2E)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK = RGBColor(0xEC, 0x48, 0x99)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def dbg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

def hero(title, sub="", accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    al = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.6), Inches(1.2), Pt(4))
    al.fill.solid(); al.fill.fore_color.rgb = accent; al.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(2.5))
    t.text_frame.word_wrap = True
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(52); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    if sub:
        p2 = t.text_frame.add_paragraph()
        p2.text = sub; p2.font.size = Pt(22); p2.font.color.rgb = MUTED; p2.space_before = Pt(20)

def sec(num, title, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    nt = s.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(3), Inches(3))
    nt.text_frame.paragraphs[0].text = f"{num:02d}"
    nt.text_frame.paragraphs[0].font.size = Pt(120); nt.text_frame.paragraphs[0].font.bold = True
    nt.text_frame.paragraphs[0].font.color.rgb = accent
    tt = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1.5))
    tt.text_frame.paragraphs[0].text = title
    tt.text_frame.paragraphs[0].font.size = Pt(40); tt.text_frame.paragraphs[0].font.bold = True
    tt.text_frame.paragraphs[0].font.color.rgb = WHITE

def stat(title, cards):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(28); t.text_frame.paragraphs[0].font.color.rgb = MUTED
    n = len(cards); cw = min(Inches(2.8), Inches(10.5) / n); gap = Inches(0.4)
    tw = n * cw + (n - 1) * gap; sx = (prs.slide_width - int(tw)) / 2
    for i, (val, lbl, clr) in enumerate(cards):
        x = int(sx) + i * (int(cw) + int(gap))
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), int(cw), Inches(3.5))
        bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
        ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), int(cw), Pt(4))
        ab.fill.solid(); ab.fill.fore_color.rgb = clr; ab.line.fill.background()
        vt = s.shapes.add_textbox(x + Inches(0.3), Inches(2.8), int(cw) - Inches(0.6), Inches(1.5))
        vt.text_frame.paragraphs[0].text = val
        vt.text_frame.paragraphs[0].font.size = Pt(48); vt.text_frame.paragraphs[0].font.bold = True
        vt.text_frame.paragraphs[0].font.color.rgb = clr
        lt = s.shapes.add_textbox(x + Inches(0.3), Inches(4.3), int(cw) - Inches(0.6), Inches(1))
        lt.text_frame.word_wrap = True
        lt.text_frame.paragraphs[0].text = lbl
        lt.text_frame.paragraphs[0].font.size = Pt(18); lt.text_frame.paragraphs[0].font.color.rgb = LIGHT

def bullets(title, items, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    al = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(0.8), Pt(3))
    al.fill.solid(); al.fill.fore_color.rgb = accent; al.line.fill.background()
    ct = s.shapes.add_textbox(Inches(1.5), Inches(1.9), Inches(10.3), Inches(5))
    ct.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = ct.text_frame.paragraphs[0] if i == 0 else ct.text_frame.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.level = item[1]
            p.font.size = Pt(18 if item[1] > 0 else 24); p.font.color.rgb = MUTED if item[1] > 0 else LIGHT
        elif item == "":
            p.text = ""; p.font.size = Pt(10)
        else:
            p.text = item; p.font.size = Pt(24); p.font.color.rgb = LIGHT
        p.space_after = Pt(6)

def steps(title, data):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    cw = Inches(3.8); ch = Inches(1.6); mg = Inches(0.3)
    colors = [CYAN, PURPLE, PINK, GREEN, ORANGE, RED]
    for i, (num, label, desc) in enumerate(data):
        col = i % 3; row = i // 3
        x = int(Inches(0.6) + col * (cw + mg)); y = int(Inches(1.8) + row * (ch + mg))
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, int(cw), int(ch))
        bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
        ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, int(cw), Pt(3))
        ab.fill.solid(); ab.fill.fore_color.rgb = colors[i % len(colors)]; ab.line.fill.background()
        ci = s.shapes.add_shape(MSO_SHAPE.OVAL, x + int(Inches(0.15)), y + int(Inches(0.2)), int(Inches(0.45)), int(Inches(0.45)))
        ci.fill.solid(); ci.fill.fore_color.rgb = colors[i % len(colors)]; ci.line.fill.background()
        ci.text_frame.paragraphs[0].text = str(num)
        ci.text_frame.paragraphs[0].font.size = Pt(16); ci.text_frame.paragraphs[0].font.bold = True
        ci.text_frame.paragraphs[0].font.color.rgb = WHITE; ci.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ci.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        lt = s.shapes.add_textbox(x + int(Inches(0.7)), y + int(Inches(0.15)), int(cw - Inches(0.9)), int(Inches(0.4)))
        lt.text_frame.paragraphs[0].text = label
        lt.text_frame.paragraphs[0].font.size = Pt(16); lt.text_frame.paragraphs[0].font.bold = True
        lt.text_frame.paragraphs[0].font.color.rgb = WHITE
        dt = s.shapes.add_textbox(x + int(Inches(0.15)), y + int(Inches(0.7)), int(cw - Inches(0.3)), int(ch - Inches(0.75)))
        dt.text_frame.word_wrap = True
        dt.text_frame.paragraphs[0].text = desc
        dt.text_frame.paragraphs[0].font.size = Pt(14); dt.text_frame.paragraphs[0].font.color.rgb = MUTED


def flow(title, steps, colors=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    active = [st for st in steps if st]
    n = len(active); pw = Inches(1.5); ph = Inches(1.1); gap = Inches(0.15)
    total = n * pw + (n - 1) * gap; sx = (prs.slide_width - int(total)) / 2; y = Inches(3.2)
    dc = colors or [CYAN, PURPLE, PINK, GREEN, ORANGE, CYAN, PURPLE]
    for i, step in enumerate(active):
        x = int(sx) + i * (int(pw) + int(gap))
        if i > 0:
            conn = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - int(gap), int(y) + int(ph // 2) - Pt(2), int(gap), Pt(4))
            conn.fill.solid(); conn.fill.fore_color.rgb = RGBColor(0x33, 0x40, 0x55); conn.line.fill.background()
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, int(y), int(pw), int(ph))
        pill.fill.solid(); pill.fill.fore_color.rgb = dc[i % len(dc)]; pill.line.fill.background()
        pill.text_frame.word_wrap = True; pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        pill.text_frame.paragraphs[0].text = step
        pill.text_frame.paragraphs[0].font.size = Pt(13); pill.text_frame.paragraphs[0].font.bold = True
        pill.text_frame.paragraphs[0].font.color.rgb = WHITE
        pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def tbl(title, headers, rows, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    cols = len(headers); n = len(rows) + 1
    ts = s.shapes.add_table(n, cols, Inches(1.2), Inches(1.7), Inches(10.9), Inches(0.45) * n)
    table = ts.table
    for j, h in enumerate(headers):
        c = table.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = SURFACE
        for par in c.text_frame.paragraphs:
            par.font.size = Pt(16); par.font.bold = True; par.font.color.rgb = accent
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = table.cell(i + 1, j); c.text = str(v)
            c.fill.solid(); c.fill.fore_color.rgb = CARD if i % 2 == 0 else BG
            for par in c.text_frame.paragraphs:
                par.font.size = Pt(15); par.font.color.rgb = LIGHT


# ==================== SLIDES ====================

hero("購買管理・出張管理\n想定運用ガイド", "next-procurement-poc  |  2026-03-28  v2")

# 01
sec(1, "システム概要")

stat("3つの役割", [
    ("申請者", "申請・発注・検収・証憑添付\n全件自分で発注", CYAN),
    ("承認者", "部門長\nDMでボタン1クリック承認", PURPLE),
    ("管理本部", "仕訳・照合・支払\n経理処理を専任", PINK),
])

bullets("1日の流れ", [
    "申請者",
    ("/purchase → 承認待ち → カード決済 → 検収 → 証憑PDF添付", 1),
    ("/mystatus で自分の未対応案件を確認 / マイページで証憑UP", 1),
    "",
    "承認者（部門長）",
    ("Bot DM の [承認] / [差戻し] — 10秒/件 / 24時間後に自動リマインド", 1),
    "",
    "管理本部",
    ("09:00 日次サマリ → 証憑確認 → 仕訳登録 → 週次で突合レポート", 1),
])

# 02
sec(2, "購買パターン")

flow("パターンA: カード払い", [
    "/purchase\n申請", "部門長\n承認", "申請者\nカード決済", "検収完了", "証憑添付", "OCR\n金額照合", "仕訳登録"])

flow("パターンB: 請求書払い", [
    "/purchase\n申請", "部門長\n承認", "申請者\n発注", "検収完了", "請求書\n提出", "管理本部\n支払処理", "仕訳登録"])

bullets("パターンC: 購入済（立替精算）", [
    "緊急時のみ: すでに購入済みの場合",
    ("申請時に証憑添付が必須", 1),
    ("承認→発注・検収スキップ→即「証憑完了」", 1),
    ("MF経費経由で給与精算", 1),
    "",
    "提出方法（全パターン共通）:",
    ("Slackスレッドにドラッグ&ドロップ", 1),
    ("マイページ（/purchase/my）からアップロード", 1),
], accent=PINK)

# 03
sec(3, "申請チャネル", CYAN)

bullets("2つの申請方法", [
    "Slackモーダル（/purchase → 「Slackモーダルで入力」）",
    ("シンプルな申請向き / 7項目入力で完了", 1),
    "",
    "Webフォーム（/purchase → 「Webフォームで入力」）",
    ("4ステップ入力 / URL自動解析 / 承認ルートプレビュー", 1),
    ("重複チェック・勘定科目推定（確認画面）", 1),
    ("一括申請（複数品目）/ 下書き自動保存 / 過去申請複製", 1),
    "",
    "ブックマークレット（/bookmarklet で初回設定）",
    ("ECサイトの商品ページからワンクリックでフォーム起動", 1),
    ("Amazon/モノタロウ/ASKUL/ヨドバシ/ビックカメラ対応", 1),
], accent=CYAN)

# 04
sec(4, "承認ルールと権限")

stat("承認と発注", [
    ("承認", "部門長のみ\n金額に関わらず1段階", GREEN),
    ("発注", "申請者が全件発注\nカード/請求書問わず", CYAN),
    ("経理", "管理本部\n仕訳・照合・支払", PINK),
])

tbl("ボタン権限マトリクス", ["操作", "申請者", "部門長", "管理本部"],
    [
        ["承認 / 差戻し", "-", "○", "-"],
        ["発注完了", "○", "○", "○"],
        ["検収完了", "○", "-", "-"],
        ["取消", "○（発注前のみ）", "-", "○"],
        ["仕訳登録", "-", "-", "○"],
    ])

# 05
sec(5, "証憑管理\n3層アーキテクチャ", PINK)

stat("証憑管理の3層", [
    ("Block", "第1層: システムブロック\n証憑なし = 仕訳不可", RED),
    ("Remind", "第2層: 段階催促\nDay 0/1/3/7 自動リマインド", ORANGE),
    ("Deter", "第3層: 抑止統制\n公開投稿・部門長エスカレ", PURPLE),
])

tbl("段階催促タイムライン", ["経過", "方法", "対象", "効果"],
    [
        ["Day 0", "スレッド投稿", "申請者", "初回通知"],
        ["Day 1", "DMダイジェスト（まとめて1通）", "申請者", "個別リマインド"],
        ["Day 3", "公開投稿（@メンション）", "全員に可視", "同僚の目"],
        ["Day 7", "部門長DMエスカレ", "部門長", "管理圧力"],
        ["承認時", "未提出一覧を表示", "部門長", "承認判断材料"],
    ], accent=ORANGE)

# 06
sec(6, "統制設計")

bullets("3層統制フレームワーク", [
    "事前統制（Preventive）",
    ("部門長承認 / 高額品は用途・理由必須 / 購入済は証憑必須", 1),
    "",
    "発見統制（Detective）",
    ("カード明細突合（週次）→ 未申請購入・承認前購入・金額不一致", 1),
    ("OCR金額照合（証憑添付時）→ 許容差 ±500円", 1),
    ("日次サマリ（毎朝09:00）→ 要対応/フォロー要/順調", 1),
    "",
    "抑止統制（Behavioral Change）",
    ("段階催促 Day 0/1/3/7 / 突合可視化 / 月次レビュー", 1),
], accent=GREEN)

# 07
sec(7, "MF連携 + 出張管理")

flow("MF連携フロー", [
    "MFビジネス\nカード", "MFクラウド\n経費", "購買台帳\n(GAS)", "OCR\n金額照合", "MF会計\nPlus", "仕訳\n完了"],
    colors=[CYAN, PURPLE, GREEN, ORANGE, PINK, GREEN])

flow("/trip 出張申請フロー", [
    "/trip\n申請", "#出張CH\n投稿", "じゃらん/\nANA/JAL予約", "MFカード\n決済", "CSV取込\n(月次)", "MF経費\n→仕訳"],
    colors=[ORANGE, ORANGE, PINK, CYAN, GREEN, GREEN])

# 08
sec(8, "Web機能")

bullets("マイページ + /mystatus", [
    "/mystatus（Slack）",
    ("自分の未対応案件をDMでサマリ表示 + マイページリンク", 1),
    "",
    "マイページ（/purchase/my）",
    ("未対応事項ダッシュボード — 発注未完了・証憑待ちを優先表示", 1),
    ("証憑待ち案件から [証憑UP] ボタンでWeb提出", 1),
    ("サマリカード + フィルター + Slackリンク", 1),
    "",
    "購買ダッシュボード（/dashboard）",
    ("ステータス分布 / 部門別 / 購入先TOP", 1),
], accent=GREEN)

# 09
sec(9, "カード明細照合", ORANGE)

steps("月次カード明細照合フロー", [
    (1, "CSV準備", "MFビジネスカード管理画面\nから利用明細CSVを\nダウンロード"),
    (2, "照合UIにアップロード", "/admin/card-matching を開き\n月を選択してCSVを\nドロップ"),
    (3, "自動照合", "予測テーブル照合(Phase1)\n→ スコア照合(Phase2)\n→ 4区分に振り分け"),
    (4, "要確認の処理", "差異タグを見て\n正しい明細を選択\n→ [これに確定]"),
    (5, "未申請利用の対応", "本人に確認通知\nまたは経費として処理"),
    (6, "引落照合", "入出金履歴CSVと\n未払金合計を突合\n→ 差額調査"),
])

tbl("照合結果の4タブ", ["タブ", "内容", "対応"],
    [
        ["自動照合済み", "申請と明細が自動マッチ", "確認のみ（差額は調整仕訳済み）"],
        ["要確認", "候補が複数 or スコア中程度", "[これに確定] で正しい明細を選択"],
        ["明細なし", "申請はあるが明細にない", "翌月繰越 or キャンセル確認"],
        ["未申請利用", "明細あるが申請がない", "本人に確認 or 経費処理"],
    ], accent=ORANGE)

bullets("マッチングの仕組み", [
    "Phase 1: 予測マッチング（高精度）",
    ("購買承認時に card_last4 × 金額 × 予想日付 を予測テーブルに自動記録", 1),
    ("金額完全一致 → 自動確定、5%以内 → 自動（差額調整付き）、5-10% → 要確認", 1),
    "",
    "Phase 2: スコアリング（フォールバック）",
    ("金額(50点) + 日付(30点) + 加盟店名(20点) = 100点満点", 1),
    ("80点以上 → 自動確定、50-79点 → 要確認、50点未満 → 未マッチ", 1),
    "",
    "引落照合:",
    ("MF会計Plusの未払金(請求)仕訳をカード別に集計 → CSVの引落額と突合", 1),
    ("差額がある場合: 月末繰越・返品・年会費等のガイドを表示", 1),
], accent=ORANGE)

# 10
sec(10, "定期タスク", GREEN)

tbl("オペレーションカレンダー", ["頻度", "時刻", "実行者", "内容"],
    [
        ["日次", "09:00", "Bot", "日次サマリ（#purchase-ops）"],
        ["日次", "10:00", "Bot", "証憑催促DM + 承認リマインド + 発注完了リマインド"],
        ["日次", "随時", "部門長", "DM承認依頼の処理"],
        ["日次", "随時", "管理本部", "証憑確認・仕訳待ち処理"],
        ["週次", "月曜 11:00", "Bot", "カード明細突合バッチ → Slack通知"],
        ["月次", "月初", "管理本部", "カード明細照合UI → 4タブ処理 + 引落照合"],
        ["月次", "~15日", "管理本部", "じゃらんCSV取込"],
        ["月次", "月末", "管理本部", "仕訳一括・コンプライアンスレビュー"],
    ], accent=GREEN)

# 11
sec(11, "ステータス遷移")

flow("通常フロー（カード・請求書共通）", [
    "申請済", "承認済", "発注済", "検収済", "証憑完了", "計上済", "支払済"])

flow("購入済フロー（立替）", [
    "申請済\n+証憑", "承認済", "検収済\n（即時）", "証憑完了", "計上済", "支払済"],
    colors=[PINK, PURPLE, GREEN, ORANGE, GREEN, GREEN])

# End
hero("Ready to Launch.", "全機能完成 → 内部テスト → 旧WFから一括移行", PURPLE)

out = "C:/Users/takeshi.izawa/.claude/projects/next-procurement-poc/docs/operational-guide.pptx"
prs.save(out)
print(f"Saved: {out}")

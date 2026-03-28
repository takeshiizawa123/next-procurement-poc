"""利用者向け運用マニュアル — Silicon Valley Style PPT (Final)"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG_DIR = os.path.join(os.path.dirname(__file__), "..", "images")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0F, 0x0F, 0x1A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK = RGBColor(0xEC, 0x48, 0x99)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SURFACE = RGBColor(0x1A, 0x1C, 0x2E)


def dbg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG


def hero(title, sub="", accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    al = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.6), Inches(1.2), Pt(4))
    al.fill.solid(); al.fill.fore_color.rgb = accent; al.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(2.5))
    t.text_frame.word_wrap = True
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(52); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    if sub:
        p2 = t.text_frame.add_paragraph()
        p2.text = sub; p2.font.size = Pt(22); p2.font.color.rgb = MUTED; p2.space_before = Pt(20)


def sec(num, title, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    nt = s.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(3), Inches(3))
    nt.text_frame.paragraphs[0].text = f"{num:02d}"
    nt.text_frame.paragraphs[0].font.size = Pt(120); nt.text_frame.paragraphs[0].font.bold = True
    nt.text_frame.paragraphs[0].font.color.rgb = accent
    tt = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(1.5))
    tt.text_frame.paragraphs[0].text = title
    tt.text_frame.paragraphs[0].font.size = Pt(40); tt.text_frame.paragraphs[0].font.bold = True
    tt.text_frame.paragraphs[0].font.color.rgb = WHITE


def bullets(title, items, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    al = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(0.8), Pt(3))
    al.fill.solid(); al.fill.fore_color.rgb = accent; al.line.fill.background()
    ct = s.shapes.add_textbox(Inches(1.5), Inches(1.9), Inches(10.3), Inches(5))
    ct.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = ct.text_frame.paragraphs[0] if i == 0 else ct.text_frame.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.level = item[1]
            p.font.size = Pt(15 if item[1] > 0 else 20); p.font.color.rgb = MUTED if item[1] > 0 else LIGHT
        elif item == "":
            p.text = ""; p.font.size = Pt(8)
        else:
            p.text = item; p.font.size = Pt(20); p.font.color.rgb = LIGHT
        p.space_after = Pt(4)


def steps(title, data):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    cw = Inches(3.8); ch = Inches(1.6); mg = Inches(0.3)
    colors = [CYAN, PURPLE, PINK, GREEN, ORANGE, RED]
    for i, (num, label, desc) in enumerate(data):
        col = i % 3; row = i // 3
        x = int(Inches(0.6) + col * (cw + mg)); y = int(Inches(1.8) + row * (ch + mg))
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, int(cw), int(ch))
        bx.fill.solid(); bx.fill.fore_color.rgb = CARD; bx.line.fill.background()
        ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, int(cw), Pt(3))
        ab.fill.solid(); ab.fill.fore_color.rgb = colors[i % len(colors)]; ab.line.fill.background()
        ci = s.shapes.add_shape(MSO_SHAPE.OVAL, x + int(Inches(0.15)), y + int(Inches(0.2)), int(Inches(0.45)), int(Inches(0.45)))
        ci.fill.solid(); ci.fill.fore_color.rgb = colors[i % len(colors)]; ci.line.fill.background()
        ci.text_frame.paragraphs[0].text = str(num)
        ci.text_frame.paragraphs[0].font.size = Pt(16); ci.text_frame.paragraphs[0].font.bold = True
        ci.text_frame.paragraphs[0].font.color.rgb = WHITE; ci.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ci.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        lt = s.shapes.add_textbox(x + int(Inches(0.7)), y + int(Inches(0.15)), int(cw - Inches(0.9)), int(Inches(0.4)))
        lt.text_frame.paragraphs[0].text = label
        lt.text_frame.paragraphs[0].font.size = Pt(16); lt.text_frame.paragraphs[0].font.bold = True
        lt.text_frame.paragraphs[0].font.color.rgb = WHITE
        dt = s.shapes.add_textbox(x + int(Inches(0.15)), y + int(Inches(0.7)), int(cw - Inches(0.3)), int(ch - Inches(0.75)))
        dt.text_frame.word_wrap = True
        dt.text_frame.paragraphs[0].text = desc
        dt.text_frame.paragraphs[0].font.size = Pt(12); dt.text_frame.paragraphs[0].font.color.rgb = MUTED


def screenshot(title, img_name, caption=""):
    """スクリーンショットを1スライドに埋め込む"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.8))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(24); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        # 画像を中央配置（最大幅11in、最大高さ5.5in）
        from PIL import Image
        with Image.open(img_path) as im:
            w, h = im.size
        max_w, max_h = Inches(11), Inches(5.5)
        ratio = min(max_w / w, max_h / h)
        iw, ih = int(w * ratio), int(h * ratio)
        x = (prs.slide_width - iw) // 2
        s.shapes.add_picture(img_path, x, Inches(1.2), iw, ih)
    if caption:
        ct = s.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11), Inches(0.5))
        ct.text_frame.paragraphs[0].text = caption
        ct.text_frame.paragraphs[0].font.size = Pt(14); ct.text_frame.paragraphs[0].font.color.rgb = MUTED


def stat(title, cards):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(28); t.text_frame.paragraphs[0].font.color.rgb = MUTED
    n = len(cards); cw = min(Inches(2.8), Inches(10.5) / n); gap = Inches(0.4)
    tw = n * cw + (n - 1) * gap; sx = (prs.slide_width - int(tw)) / 2
    for i, (val, lbl, clr) in enumerate(cards):
        x = int(sx) + i * (int(cw) + int(gap))
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), int(cw), Inches(3.5))
        bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B); bx.line.fill.background()
        ab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), int(cw), Pt(4))
        ab.fill.solid(); ab.fill.fore_color.rgb = clr; ab.line.fill.background()
        vt = s.shapes.add_textbox(x + Inches(0.3), Inches(2.8), int(cw) - Inches(0.6), Inches(1.5))
        vt.text_frame.paragraphs[0].text = val
        vt.text_frame.paragraphs[0].font.size = Pt(48); vt.text_frame.paragraphs[0].font.bold = True
        vt.text_frame.paragraphs[0].font.color.rgb = clr
        lt2 = s.shapes.add_textbox(x + Inches(0.3), Inches(4.3), int(cw) - Inches(0.6), Inches(1))
        lt2.text_frame.word_wrap = True
        lt2.text_frame.paragraphs[0].text = lbl
        lt2.text_frame.paragraphs[0].font.size = Pt(16); lt2.text_frame.paragraphs[0].font.color.rgb = LIGHT


def tbl(title, headers, rows, accent=CYAN):
    s = prs.slides.add_slide(prs.slide_layouts[6]); dbg(s)
    t = s.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(10), Inches(1))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(32); t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    cols = len(headers); n = len(rows) + 1
    ts = s.shapes.add_table(n, cols, Inches(1.2), Inches(1.7), Inches(10.9), Inches(0.45) * n)
    table = ts.table
    for j, h in enumerate(headers):
        c = table.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = SURFACE
        for par in c.text_frame.paragraphs:
            par.font.size = Pt(13); par.font.bold = True; par.font.color.rgb = accent
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = table.cell(i + 1, j); c.text = str(v)
            c.fill.solid(); c.fill.fore_color.rgb = CARD if i % 2 == 0 else BG
            for par in c.text_frame.paragraphs:
                par.font.size = Pt(12); par.font.color.rgb = LIGHT


# ==================== SLIDES ====================

# ─────────────────────────────────────────────
# Part 0: タイトル + 全体フロー
# ─────────────────────────────────────────────

hero("購買申請システム\n利用者マニュアル", "Slack Bot + Webフォーム で購買申請から仕訳まで\nv1.0 — 2026-03-28")

tbl("このマニュアルで使う用語", ["用語", "意味"],
    [
        ["証憑（しょうひょう）", "購入を証明する書類。納品書・領収書・請求書など"],
        ["検収（けんしゅう）", "届いた物品が注文通りか確認すること"],
        ["仕訳（しわけ）", "会計処理。管理本部が行うので申請者は意識不要"],
        ["PO番号", "申請ごとに自動発行される管理番号（例: PO-202603-0050）"],
    ], accent=MUTED)

steps("全体フロー — 購買申請から経理処理まで", [
    (1, "申請", "申請者\n/purchase で入力"),
    (2, "承認", "部門長\nボタン1つで承認"),
    (3, "発注", "申請者 or 管理本部\nMFカードで購入"),
    (4, "検収", "申請者\n届いたらボタン1つ"),
    (5, "証憑", "申請者\n納品書をスレッドに添付"),
    (6, "仕訳・照合・引落", "管理本部\nMF会計Plus登録\nカード明細突合"),
])

stat("3つの役割 — 誰が何をする？",
    [("申請者", "①申請 → ③発注 → ④検収 → ⑤証憑\n4ステップだけ", CYAN),
     ("部門長", "②承認\nボタン1つ（10秒）", PURPLE),
     ("管理本部", "⑥仕訳 → ⑦照合 → ⑧引落\n経理処理を一手に担当", PINK)])

# ─────────────────────────────────────────────
# Part A: 申請者の目線
# ─────────────────────────────────────────────

hero("Part A\n申請者ガイド", "①申請 → ③発注 → ④検収 → ⑤証憑", CYAN)

# A-1: 申請
sec(1, "申請する", CYAN)

steps("Slackモーダル or Webフォームで申請", [
    (1, "/purchase を入力", "Slackの任意チャンネルで\nコマンドを送信"),
    (2, "入力方法を選択", "Slackモーダル（簡易）\nor Webフォーム（高機能）"),
    (3, "フォームに入力", "品目・金額・数量・支払方法\n・購入目的を入力"),
    (4, "送信", "#purchase-request に\n自動投稿されます"),
    (5, "承認を待つ", "部門長にDMで\n承認依頼が届きます"),
    (6, "番号を確認", "PO-202603-0050 のような\n申請番号が発行されます"),
])

tbl("入力項目", ["項目", "必須", "説明", "入力例"],
    [
        ["申請区分", "★", "「購入前」or「購入済」", "購入前"],
        ["品目名", "★", "何を買うか", "会議用モニター"],
        ["単価（税抜）", "★", "1個の金額", "45,000"],
        ["数量", "★", "個数", "1"],
        ["支払方法", "★", "カード or 請求書", "会社カード"],
        ["購入先", "★", "店舗/サイト名", "Amazon"],
        ["購入目的", "★", "利用目的", "業務利用"],
        ["商品URL", "", "自動で情報取得", "https://..."],
    ])

# A-2: 承認後 → 発注
sec(2, "承認後: 発注する", CYAN)

bullets("承認された後の操作", [
    "会社カード + 10万円未満 → 自分で発注",
    ("DM: 「承認されました。発注してください」", 1),
    ("MFバーチャルカードで購入", 1),
    ("購入後 → #purchase-request の [発注完了] ボタンを押す", 1),
    ("3日後に自動リマインドDM（押し忘れ防止）", 1),
    "",
    "会社カード + 10万円以上 / 請求書払い",
    ("管理本部が代行するので、あなたの操作は不要です", 1),
    "",
    "差戻しの場合 → DMで理由が届く → /purchase で再申請",
    "取消 → [取消] ボタン（発注前のみ）",
])

# A-3: 検収 → 証憑
sec(3, "検収と証憑の提出", CYAN)

steps("検収 → 証憑提出 → 完了", [
    (1, "物品の到着確認", "注文通りの内容か\n確認してください"),
    (2, "[検収完了] を押す", "#purchase-request の\n該当メッセージで実行"),
    (3, "証憑を用意", "納品書 or 領収書の\nPDF / 写真を準備"),
    (4, "スレッドに添付", "該当申請のスレッドに\nドラッグ&ドロップ"),
    (5, "Bot が自動検知", "種別判定 + OCR金額照合\n→ 結果がスレッドに表示"),
    (6, "完了!", "「あなたの作業は完了です」\n経理処理は管理本部が対応"),
])

bullets("証憑を出さないとどうなる？", [
    "証憑がないと仕訳・支払に進めません",
    "",
    "自動リマインド:",
    ("翌日 — DMで「証憑を提出してください」", 1),
    ("3日後 — スレッドに公開投稿（@メンション）", 1),
    ("7日後 — 部門長にエスカレーション", 1),
    "",
    "次回の申請時、部門長の承認画面に「未提出一覧」が表示されます",
], accent=ORANGE)

# A-4: 購入済（立替）
sec(4, "購入済申請（立替精算）", RED)

bullets("立替精算の手順", [
    "/purchase → 申請区分: 「購入済」を選択",
    ("品目・金額を入力 + レシートを必ず添付", 1),
    ("承認 → 発注・検収スキップ → 「証憑完了」", 1),
    "",
    "注意:",
    ("購入済は緊急時の例外措置です", 1),
    ("原則: 事前に /purchase → 承認後にMFカードで購入", 1),
    ("証憑なしの購入済申請は処理されません", 1),
], accent=RED)

bullets("立替精算とMF経費の関係", [
    "立替精算の経理処理はMF経費を経由します",
    "",
    "処理フロー:",
    ("① 申請者が /purchase「購入済」で申請 + 証憑添付", 1),
    ("② 部門長が承認", 1),
    ("③ システムがMF経費に証憑をアップロード（自動）", 1),
    ("④ 管理本部がMF経費で経費精算を確定", 1),
    ("⑤ 給与と合算して立替分を振込", 1),
    "",
    "申請者がMF経費を直接操作する必要はありません",
    ("MF経費のカード明細から経費登録しないでください（二重計上防止）", 1),
], accent=RED)

# A-5: 出張
sec(5, "出張申請（/trip）", ORANGE)

steps("出張申請フロー", [
    (1, "/trip を入力", "Slackのどのチャンネル\n・DMからでも実行可能"),
    (2, "モーダルに入力", "行き先・日程・目的\n交通手段・概算額・宿泊先"),
    (3, "送信", "#出張チャンネルに\n自動投稿されます"),
    (4, "予約・手配", "各交通機関・宿泊を\n予約（次スライド参照）"),
    (5, "MFカードで決済", "すべてMFビジネス\nカードで支払い"),
    (6, "精算は自動", "カード明細→MF経費\n→管理本部が仕訳"),
])

tbl("/trip 入力項目", ["項目", "必須", "入力例"],
    [
        ["出張先", "★", "大阪本社"],
        ["出発日", "★", "2026-04-01"],
        ["帰着日", "★", "2026-04-03"],
        ["出張目的", "★", "クライアント打合せ"],
        ["利用交通手段", "★", "新幹線のぞみ 東京→新大阪"],
        ["概算額（円）", "★", "45000"],
        ["宿泊先", "", "じゃらんで予約済み / ホテル名"],
    ], accent=ORANGE)

tbl("交通手段別の予約・決済方法", ["手段", "予約方法", "決済", "備考"],
    [
        ["新幹線", "スマートEX / えきねっと", "MFカード", "EX予約の場合はICカードで乗車"],
        ["航空券（ANA）", "ANA Biz（法人サイト）", "MFカード", "法人IDでログイン"],
        ["航空券（JAL）", "JAL Online（法人サイト）", "MFカード", "法人IDでログイン"],
        ["レンタカー", "各社Webサイト", "MFカード", "利用交通手段欄に「レンタカー」と記入"],
        ["タイムズカー", "タイムズカーアプリ", "MFカード", "法人カードを事前登録"],
        ["タクシー", "現地で利用", "MFカード", "領収書不要（カード明細で確認）"],
    ], accent=ORANGE)

bullets("宿泊の手配", [
    "じゃらんJCS（法人向けサービス）を使用",
    ("法人専用項目1にプロジェクト番号を入力", 1),
    ("支払は会社宛ての請求書払い（個人負担なし）", 1),
    "",
    "楽天トラベルRacco も今後利用可能になります",
    "",
    "出張後の精算:",
    ("カード決済分 → MF経費に自動連携（操作不要）", 1),
    ("宿泊費 → じゃらんCSV一括取込（管理本部が月次処理）", 1),
    ("立替が発生した場合 → /purchase「購入済」で申請", 1),
], accent=ORANGE)

# A-6: マイページ
sec(6, "マイページ", GREEN)

bullets("マイページでできること", [
    "アクセス: https://{システムURL}/purchase/my",
    "",
    "申請一覧の確認",
    ("自分が出した全申請のステータスを一覧表示", 1),
    ("ステータスバッジ: 承認済 / 発注済 / 証憑待ち / 計上済 etc.", 1),
    ("各申請からSlackスレッドへワンクリックでジャンプ", 1),
    "",
    "未対応アラート（黄色い帯）",
    ("証憑未提出の案件が「○件あります」と目立つ表示", 1),
    ("経過日数が赤字で表示される（3日経過、5日経過...）", 1),
    "",
    "証憑アップロード",
    ("証憑待ち案件をクリック → 右パネルのドロップエリアにファイルを投入", 1),
    ("Slackが使えない状況でもWebから証憑提出が可能", 1),
], accent=GREEN)

screenshot("マイページ画面", "mypage.png", "申請一覧 + 証憑未提出アラート + 詳細パネル")

# A-7: ブックマークレット
sec(7, "ブックマークレット", GREEN)

steps("ECサイトからワンクリック申請", [
    (1, "初回設定（1回のみ）", "/bookmarklet ページを開く\n→「購買申請」ボタンを\nブックマークバーにD&D"),
    (2, "商品ページを開く", "Amazon・モノタロウ等で\n購入したい商品の\nページを表示"),
    (3, "ブックマークをクリック", "ブックマークバーの\n「購買申請」を押す"),
    (4, "フォームが自動入力済み", "商品名・価格・購入先\n・URLが入力された\n状態で開きます"),
    (5, "残りを入力して送信", "支払方法・数量・目的\nを追加入力して\n申請完了"),
])

tbl("URL自動解析 対応サイト", ["サイト", "商品名", "価格", "備考"],
    [
        ["Amazon.co.jp", "○", "○", "サーバーブロック時はURL文字列から商品名を推定"],
        ["モノタロウ", "○", "○", "安定して取得可能"],
        ["ASKUL", "○", "○", "安定して取得可能"],
        ["ヨドバシ.com", "○", "○", ""],
        ["ビックカメラ", "○", "○", ""],
        ["その他サイト", "ページタイトル", "—", "ブックマークレット経由のみ"],
    ], accent=GREEN)

# ─────────────────────────────────────────────
# Part B: 承認者（部門長）の目線
# ─────────────────────────────────────────────

hero("Part B\n承認者（部門長）ガイド", "②承認 — ボタン1つで完了", PURPLE)

sec(8, "承認する", PURPLE)

bullets("承認者の操作ガイド", [
    "承認依頼: 部下の申請時にBot DMが届く",
    ("内容: 品目名・金額・購入先・支払方法", 1),
    "",
    "承認: [承認] ボタンを押すだけ（10秒）",
    ("DMからでも #purchase-request からでもOK", 1),
    ("10万以上は → 管理本部の二段階目に自動回付", 1),
    ("24時間放置するとリマインドDMが届きます", 1),
    "",
    "差戻し: [差戻し] → 理由入力 → 送信",
    "",
    "確認ポイント:",
    ("申請者の証憑未提出一覧が表示される場合あり", 1),
    ("未提出が多い → 先に提出を促してから承認を推奨", 1),
], accent=PURPLE)

screenshot("承認待ちメッセージ", "slack-pending.png", "#purchase-request に投稿される申請メッセージ")
screenshot("承認後の表示", "slack-approved.png", "承認すると [発注完了] [検収完了] ボタンに切り替わる")
screenshot("差戻し後の表示", "slack-rejected.png", "差戻し理由が表示される。申請者にDMで通知")

# ─────────────────────────────────────────────
# Part C: 管理本部ガイド
# ─────────────────────────────────────────────

hero("Part C\n管理本部ガイド", "⑥仕訳 → ⑦照合 → ⑧引落 + 発注代行", PINK)

sec(9, "日常業務", PINK)

bullets("管理本部の日常", [
    "毎朝 09:00: #purchase-ops 日次サマリ確認",
    ("🔴要対応（承認待ち・発注待ち）", 1),
    ("🟡フォロー要（3日以上停滞の証憑待ち）", 1),
    ("🟢順調（進行中・完了件数）", 1),
    "",
    "発注処理: [開く] → 発注 → [発注完了]",
    ("対象: 10万以上 + 請求書払いのみ", 1),
    "",
    "仕訳登録: 証憑完了案件を MF会計Plus に登録",
    ("[仕訳登録] で自動ドラフト作成 → 借方:費用科目 / 貸方:未払金 or 買掛金", 1),
], accent=PINK)

screenshot("管理ダッシュボード", "dashboard.png", "/dashboard — ステータス分布・部門別・購入先TOP")

sec(10, "カード明細照合（月次）", ORANGE)

steps("月次照合フロー", [
    (1, "CSVダウンロード", "MFビジネスカード管理画面\nから利用明細CSVを取得"),
    (2, "照合UIにアップロード", "/admin/card-matching\nCSVをドロップ → 自動照合"),
    (3, "4タブを処理", "自動照合済 / 要確認\n/ 明細なし / 未申請利用"),
    (4, "引落照合", "入出金履歴CSVと\n未払金合計を突合"),
    (5, "完了通知", "全件処理後\nSlackに完了通知"),
])

screenshot("自動照合済みタブ", "matching-confirmed.png", "差額ありは黄色で強調。差額なしは灰色で一覧表示")
screenshot("要確認タブ", "matching-review.png", "差異タグ（日付/金額/取引先名）を見て [これに確定] で選択")
screenshot("未申請利用タブ", "matching-unreported.png", "[本人に確認] でSlack通知 or [経費で処理]")

tbl("照合結果の4タブ", ["タブ", "内容", "対応"],
    [
        ["自動照合済み", "申請と明細が自動マッチ", "確認のみ（差額は調整仕訳済み）"],
        ["要確認", "候補が複数 or スコア中程度", "[これに確定] で正しい明細を選択"],
        ["明細なし", "申請はあるが明細にない", "翌月繰越 or キャンセル確認"],
        ["未申請利用", "明細あるが申請がない", "本人に確認 or 経費処理"],
    ], accent=ORANGE)

bullets("マッチングの仕組み", [
    "Phase 1: 予測マッチング — 承認時にcard_last4×金額×日付を自動記録",
    ("金額完全一致 → 自動確定、5%以内 → 差額調整、5-10% → 要確認", 1),
    "",
    "Phase 2: スコアリング — MF仕訳と明細を突合",
    ("金額(50点) + 日付(30点) + 加盟店名(20点) = 100点満点", 1),
    "",
    "引落照合: 未払金(請求)合計 vs 銀行引落額",
    ("差額 → 月末繰越/返品/年会費等のガイド表示", 1),
], accent=ORANGE)

tbl("月次作業カレンダー", ["作業", "内容", "所要時間"],
    [
        ["カード明細照合", "照合UI 4タブ処理", "30-60分"],
        ["引落照合", "引落照合タブで未払金突合", "15分"],
        ["じゃらんCSV取込", "管理画面からCSV → アップロード", "15分"],
        ["仕訳一括処理", "未処理の仕訳待ちを一括登録", "30分"],
        ["コンプライアンスレビュー", "未申請/承認前購入の件数確認", "15分"],
    ], accent=PINK)

# ─────────────────────────────────────────────
# Part D: FAQ & 早見表
# ─────────────────────────────────────────────

hero("FAQ &\n操作早見表", "", GREEN)

tbl("よくある質問", ["質問", "回答"],
    [
        ["申請を間違えた", "[取消] → 再申請（発注後は管理本部に連絡）"],
        ["証憑を紛失した", "管理本部に相談（カード明細スクショで代替可）"],
        ["承認が来ない", "部門長にSlackで確認（24時間後に自動リマインド）"],
        ["10万以上を自分で買いたい", "不可（管理本部が発注 — 固定資産管理ルール）"],
        ["カード明細照合って何？", "管理本部の月次作業。一般社員は意識不要"],
        ["未申請利用の連絡が来た", "/purchase で事後申請を出してください"],
    ], accent=GREEN)

tbl("操作早見表", ["やりたいこと", "操作", "対象"],
    [
        ["購買申請", "/purchase → モーダル or Webフォーム", "申請者"],
        ["出張申請", "/trip → モーダル", "申請者"],
        ["自分の状況確認", "/purchase/my（マイページ）", "申請者"],
        ["承認する", "DM or #purchase-request の [承認]", "部門長"],
        ["発注完了を報告", "[発注完了] ボタン", "申請者/管理"],
        ["検収完了を報告", "[検収完了] ボタン", "申請者"],
        ["証憑を提出", "スレッドにD&D or マイページ [証憑UP]", "申請者"],
        ["カード明細を照合", "/admin/card-matching → CSV読込", "管理本部"],
        ["引落額を照合", "/admin/card-matching → 引落照合タブ", "管理本部"],
    ], accent=CYAN)

# --- End ---
hero("ご不明点は\n管理本部まで", "#purchase-ops チャンネル または DM", PURPLE)

out = "C:/Users/takeshi.izawa/.claude/projects/next-procurement-poc/docs/user-manual.pptx"
prs.save(out)
print(f"Saved: {out}")

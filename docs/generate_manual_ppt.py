"""利用者向け運用マニュアル PPT生成"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x2C, 0x3E, 0x50)
BLUE = RGBColor(0x29, 0x80, 0xB9)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x95, 0xA5, 0xA6)


def bg(slide, color=LIGHT_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header_bar(slide, text, color=DARK):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE


def title_slide(title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
    t = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(2))
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = t.text_frame.add_paragraph()
        p2.text = sub; p2.font.size = Pt(20); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)


def section_slide(title, icon=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BLUE)
    t = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    p = t.text_frame.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER


def bullets_slide(title, items, note=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header_bar(s, title)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.7))
    t.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = t.text_frame.paragraphs[0] if i == 0 else t.text_frame.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.level = item[1]
            p.font.size = Pt(15 if item[1] > 0 else 18)
            p.font.color.rgb = GRAY if item[1] > 0 else DARK
        else:
            p.text = item; p.font.size = Pt(18); p.font.color.rgb = DARK
        p.space_after = Pt(6)
    if note:
        p = t.text_frame.add_paragraph()
        p.text = note; p.font.size = Pt(14); p.font.color.rgb = GRAY; p.font.italic = True; p.space_before = Pt(12)


def table_slide(title, headers, rows):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header_bar(s, title)
    cols = len(headers); n = len(rows) + 1
    ts = s.shapes.add_table(n, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.45 * n))
    tbl = ts.table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = BLUE
        for par in c.text_frame.paragraphs: par.font.size = Pt(14); par.font.bold = True; par.font.color.rgb = WHITE; par.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(v)
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
            for par in c.text_frame.paragraphs: par.font.size = Pt(13); par.font.color.rgb = DARK


def step_slide(title, steps):
    """Numbered step cards"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header_bar(s, title)
    card_w = Inches(3.8); card_h = Inches(1.6); margin = Inches(0.3)
    cols_per_row = 3
    for i, (num, label, desc) in enumerate(steps):
        col = i % cols_per_row; row = i // cols_per_row
        x = Inches(0.6) + col * (card_w + margin)
        y = Inches(1.5) + row * (card_h + margin)
        # card bg
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(card_w), int(card_h))
        bx.fill.solid(); bx.fill.fore_color.rgb = WHITE
        bx.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); bx.line.width = Pt(1)
        # number circle
        cx = int(x + Inches(0.15)); cy = int(y + Inches(0.15))
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, int(Inches(0.5)), int(Inches(0.5)))
        circle.fill.solid(); circle.fill.fore_color.rgb = BLUE; circle.line.fill.background()
        circle.text_frame.paragraphs[0].text = str(num)
        circle.text_frame.paragraphs[0].font.size = Pt(18); circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE; circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # label
        lt = s.shapes.add_textbox(int(x + Inches(0.8)), int(y + Inches(0.1)), int(card_w - Inches(1)), int(Inches(0.4)))
        lt.text_frame.paragraphs[0].text = label
        lt.text_frame.paragraphs[0].font.size = Pt(16); lt.text_frame.paragraphs[0].font.bold = True; lt.text_frame.paragraphs[0].font.color.rgb = DARK
        # desc
        dt = s.shapes.add_textbox(int(x + Inches(0.15)), int(y + Inches(0.65)), int(card_w - Inches(0.3)), int(card_h - Inches(0.7)))
        dt.text_frame.word_wrap = True
        dt.text_frame.paragraphs[0].text = desc
        dt.text_frame.paragraphs[0].font.size = Pt(12); dt.text_frame.paragraphs[0].font.color.rgb = GRAY


# ==================== SLIDES ====================

# 1
title_slide("購買申請システム\n利用者マニュアル", "Slack Bot で購買申請から仕訳計上まで\nv0.1 ドラフト — 2026-03-26")

# 2
section_slide("購買申請の出し方", "1")

# 3
step_slide("申請の手順", [
    (1, "/purchase と入力", "Slackの任意チャンネルで\n/purchase コマンドを送信"),
    (2, "入力方法を選択", "Slackモーダル（簡易）\nまたはWebフォーム（高機能）"),
    (3, "フォームに入力", "品目名・金額・数量\n支払方法・購入目的を入力"),
    (4, "送信", "#purchase-request に\n自動投稿されます"),
    (5, "承認待ち", "部門長にDMで\n承認依頼が届きます"),
    (6, "申請番号の確認", "PR-0050 のような\n番号が発行されます"),
])

# 4
table_slide("入力項目一覧", ["項目", "必須", "説明", "入力例"],
    [
        ["申請区分", "★", "「購入前」または「購入済」", "購入前"],
        ["品目名", "★", "何を購入するか", "会議用モニター"],
        ["単価（税抜）", "★", "1個あたりの金額", "45,000"],
        ["数量", "★", "購入数量", "1"],
        ["支払方法", "★", "会社カード / 請求書払い", "会社カード"],
        ["購入先", "", "購入先の名前", "Amazon"],
        ["購入目的", "★", "業務利用 / プロジェクト利用", "業務利用"],
        ["受取場所", "★", "本社 / 支社 / リモート", "本社"],
        ["商品URL", "", "ECサイトURL（自動情報取得）", "https://..."],
    ])

# 5
section_slide("承認後の操作", "2")

# 6
bullets_slide("承認された場合", [
    "会社カード + 10万円未満（自分で発注）:",
    ("Bot DMで「承認されました。発注してください」と届く", 1),
    ("MFバーチャルカードで商品を購入", 1),
    ("#purchase-request の [発注完了] ボタンを押す", 1),
    "",
    "会社カード + 10万円以上 / 請求書払い:",
    ("管理本部が発注を代行します — 申請者の操作は不要", 1),
    ("管理本部が発注完了するとDMで通知が届きます", 1),
    "",
    "差戻しされた場合:",
    ("Bot DMで理由が届く → 内容修正して /purchase で再申請", 1),
])

# 7
section_slide("検収と証憑の提出", "3")

# 8
step_slide("検収 → 証憑提出の手順", [
    (1, "物品の到着を確認", "届いた物品が\n注文通りか確認"),
    (2, "[検収完了] ボタン", "#purchase-request の\n該当メッセージで押す"),
    (3, "証憑を用意", "納品書・領収書の\nPDF or 写真を準備"),
    (4, "スレッドに添付", "該当申請のスレッドに\nファイルをドラッグ&ドロップ"),
    (5, "自動検知", "Botが種別判定 +\nOCR金額照合を実行"),
    (6, "完了", "「金額一致」表示で\n仕訳計上待ちに移行"),
])

# 9
bullets_slide("証憑を出さないとどうなる？", [
    "証憑がないと会計処理（仕訳・支払）に進めません",
    "",
    "自動リマインドの流れ:",
    ("翌日: DMで「証憑を提出してください」（まとめて1通）", 1),
    ("3日後: スレッドに公開投稿（@メンション付き — 周りにも見えます）", 1),
    ("7日後: 部門長にエスカレーション", 1),
    "",
    "さらに:",
    ("次回の申請時、部門長の承認画面に「未提出一覧」が表示されます", 1),
], note="対応形式: PDF、JPEG、PNG、HEIC、WebP、TIFF")

# 10
section_slide("購入済申請（立替精算）", "4")

# 11
bullets_slide("購入済（立替）の手順と注意事項", [
    "手順:",
    ("/purchase → 申請区分:「購入済」を選択", 1),
    ("品目・金額を入力 + 証憑（レシート）を必ず添付", 1),
    ("承認後、自動的に「検収済・証憑待ち」に遷移", 1),
    "",
    "注意事項:",
    ("購入済申請は緊急時の例外措置です", 1),
    ("原則: 事前に /purchase で申請 → MFカードで購入", 1),
    ("証憑なしの購入済申請は処理されません", 1),
])

# 12
section_slide("承認者（部門長）向け", "5")

# 13
bullets_slide("承認者の操作", [
    "承認依頼の受け取り:",
    ("部下の申請時にBotからDMが届きます", 1),
    ("内容: 品目名・金額・購入先・支払方法", 1),
    "",
    "承認する: [承認] ボタンを押すだけ（10秒）",
    ("DMからでも #purchase-request からでも操作可能", 1),
    ("10万円以上は部門長承認後に管理本部の二段階目に回ります", 1),
    "",
    "差戻しする: [差戻し] ボタン → 理由入力 → 送信",
    ("申請者にDMで理由が通知されます", 1),
    "",
    "確認ポイント:",
    ("申請者の証憑未提出一覧が表示される場合があります", 1),
    ("未提出が多い場合は先に提出を促してから承認を推奨", 1),
])

# 14
section_slide("管理本部向け", "6")

# 15
bullets_slide("管理本部の日常業務", [
    "毎朝 09:00: #purchase-ops の日次サマリを確認",
    ("🔴要対応: 承認待ち・発注待ち（管理本部が発注する案件）", 1),
    ("🟡フォロー要: 3日以上停滞の証憑待ち（スレッドリンク付き）", 1),
    ("🟢順調: 進行中・完了件数", 1),
    "",
    "発注処理: サマリの [開く] からスレッドへ → 発注 → [発注完了]",
    ("10万以上 or 請求書払いのみが管理本部の発注対象", 1),
    "",
    "仕訳登録: 証憑完了案件を MF会計Plus に登録",
    ("[仕訳登録] ボタンで自動ドラフト作成", 1),
    "",
    "週次: カード明細突合レポート確認（月曜 11:00 自動投稿）",
    "月次: じゃらんCSV取込、仕訳一括処理、コンプライアンスレビュー",
])

# 16
section_slide("出張申請（/trip）", "7")

# 17
step_slide("出張申請の手順", [
    (1, "/trip と入力", "Slackの任意チャンネルで\n/trip コマンドを送信"),
    (2, "モーダルに入力", "行き先・日程・目的\n交通手段・概算額"),
    (3, "送信", "#出張チャンネルに\n自動投稿されます"),
    (4, "予約", "じゃらん / ANA Biz /\nJAL Online で予約"),
    (5, "出張実施", "MFカードで\n交通費・食事代を決済"),
    (6, "精算（自動）", "カード明細 → MF経費\n→ 管理本部が仕訳"),
])

# 18
bullets_slide("宿泊・交通の手配", [
    "宿泊: じゃらんJCS（法人向け）を使用",
    ("予約時に「法人専用項目1」にプロジェクト番号を入力", 1),
    ("支払は会社宛て請求書払い（個人負担なし）", 1),
    ("楽天トラベルRacco も今後利用可能（準備中）", 1),
    "",
    "航空券: ANA Biz / JAL Online で法人予約",
    "",
    "新幹線: スマートEX等 → MFビジネスカードで決済",
    "",
    "その他（タクシー・食事等）: MFカードで決済",
    "",
    "立替が発生した場合: /purchase → 「購入済」で申請",
])

# 19
section_slide("よくある質問（FAQ）", "Q&A")

# 20
table_slide("FAQ", ["質問", "回答"],
    [
        ["申請を間違えた", "発注前なら [取消] → 再申請。発注後は管理本部に連絡"],
        ["証憑を紛失した", "管理本部に相談。カード明細のスクショ等で代替可能な場合あり"],
        ["承認が来ない", "部門長にSlackで直接確認してください"],
        ["10万以上を自分で買いたい", "不可。管理本部が発注・決済します（固定資産管理のルール）"],
        ["請求書払いの証憑は？", "ベンダーからの請求書をスレッドに添付"],
        ["複数品目をまとめて申請", "Webフォームなら「品目追加」で一括申請可能"],
        ["出張の承認フローは？", "/trip は投稿のみ。予約前にチャンネルで上長了承を得る運用"],
    ])

# 21
section_slide("操作早見表", "")

# 22
table_slide("操作早見表", ["やりたいこと", "操作"],
    [
        ["購買申請を出す", "/purchase → モーダルまたはWebフォーム"],
        ["出張申請を出す", "/trip → モーダル"],
        ["申請を承認する", "DM or #purchase-request の [承認] ボタン"],
        ["発注完了を報告", "#purchase-request の [発注完了] ボタン"],
        ["検収完了を報告", "#purchase-request の [検収完了] ボタン"],
        ["証憑を提出する", "申請スレッドにファイルをドラッグ&ドロップ"],
        ["申請を取り消す", "[取消] ボタン（発注前のみ）"],
    ])

# 23
title_slide("ご不明点は\n管理本部までお問い合わせください", "#purchase-ops チャンネルまたはDMで")

# Save
out = "C:/Users/takeshi.izawa/.claude/projects/next-procurement-poc/docs/user-manual-draft.pptx"
prs.save(out)
print(f"Saved: {out}")
